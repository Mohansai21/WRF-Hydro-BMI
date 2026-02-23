#!/usr/bin/env python3
"""
Generate BMI Architecture & Strategy PowerPoint Presentation
For WRF-Hydro BMI Project — Mohan Sai
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette ───
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)       # Deep navy
SLIDE_BG = RGBColor(0x0F, 0x0F, 0x23)       # Darker navy
ACCENT_BLUE = RGBColor(0x4E, 0xC9, 0xF0)    # Bright cyan
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)   # Green
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)     # Red
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)  # Orange
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)  # Purple
ACCENT_YELLOW = RGBColor(0xF1, 0xC4, 0x0F)  # Yellow
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBD, 0xBD, 0xBD)
MED_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
CARD_BG = RGBColor(0x1E, 0x2A, 0x3A)        # Darker card
CODE_BG = RGBColor(0x0D, 0x11, 0x17)        # Code block bg


def set_slide_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Round corners
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=LIGHT_GRAY, bullet_color=ACCENT_BLUE, font_name="Calibri"):
    """Add a bulleted text list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(2)
        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = font_name
        # Content
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = font_name
    return txBox


def add_code_block(slide, left, top, width, height, code_text, font_size=10):
    """Add a code block with monospace font and dark bg."""
    # Background shape
    shape = add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE, Pt(0.5))
    # Code text
    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code_text.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN
        p.font.name = "Consolas"
        p.space_after = Pt(1)
        p.space_before = Pt(1)
    return shape


def add_section_number(slide, number, top=Inches(0.3)):
    """Add a section/slide number badge."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), top, Inches(0.5), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    shape.adjustments[0] = 0.3
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)


def add_table(slide, left, top, width, rows, cols, data, col_widths=None):
    """Add a styled table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.4 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            # Style
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Calibri"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = LIGHT_GRAY
                p.alignment = PP_ALIGN.LEFT

            # Cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if r == 0:
                cell_fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
            elif r % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x35)
            else:
                cell_fill.fore_color.rgb = RGBColor(0x15, 0x1E, 0x2D)

    return table_shape


# ═══════════════════════════════════════════════════════════════
#  BUILD THE PRESENTATION
# ═══════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─────────────────────────────────────────────────
# SLIDE 1: TITLE SLIDE
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Top accent line
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

# Title
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "BMI Architecture & Strategy", font_size=44, color=WHITE, bold=True)

# Subtitle
add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
             "SCHISM vs WRF-Hydro  |  BMI Fortran Template  |  Heat Model Deep Dive",
             font_size=22, color=ACCENT_BLUE)

# Description
add_text_box(slide, Inches(1), Inches(3.8), Inches(10), Inches(1.2),
             "Understanding BMI integration patterns: why SCHISM uses inline #ifdef,\n"
             "why WRF-Hydro needs a separate wrapper, and how to build one from scratch.",
             font_size=16, color=LIGHT_GRAY)

# Info cards at bottom
add_shape(slide, Inches(1), Inches(5.5), Inches(3), Inches(0.7), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.2), Inches(5.6), Inches(2.6), Inches(0.5),
             "Project: WRF-Hydro BMI Wrapper", font_size=12, color=ACCENT_BLUE)

add_shape(slide, Inches(4.3), Inches(5.5), Inches(3), Inches(0.7), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(4.5), Inches(5.6), Inches(2.6), Inches(0.5),
             "Phase 1: Understanding & Planning", font_size=12, color=ACCENT_GREEN)

add_shape(slide, Inches(7.6), Inches(5.5), Inches(3), Inches(0.7), CARD_BG, ACCENT_ORANGE, Pt(1))
add_text_box(slide, Inches(7.8), Inches(5.6), Inches(2.6), Inches(0.5),
             "Goal: Couple WRF-Hydro + SCHISM", font_size=12, color=ACCENT_ORANGE)

# Bottom accent line
add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), ACCENT_BLUE)


# ─────────────────────────────────────────────────
# SLIDE 2: AGENDA / ROADMAP
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "Presentation Roadmap", font_size=32, color=WHITE, bold=True)

# Three section cards
sections = [
    ("Part 1", "SCHISM's BMI Approach", "How SCHISM enables BMI\nvia #ifdef USE_NWM_BMI\nin its own source code", ACCENT_BLUE, "Slides 3-6"),
    ("Part 2", "WRF-Hydro: Why Different?", "Why inline #ifdef won't work\nfor WRF-Hydro and what\napproach we'll use instead", ACCENT_ORANGE, "Slides 7-12"),
    ("Part 3", "BMI Template & Heat Model", "Deep dive into bmi.f90\nand bmi_heat.f90 — the\nblueprint for our wrapper", ACCENT_GREEN, "Slides 13-22"),
]

for i, (part, title, desc, color, slides_range) in enumerate(sections):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.5)

    # Card background
    add_shape(slide, x, y, Inches(3.7), Inches(4.5), CARD_BG, color, Pt(1.5))

    # Part number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.3), Inches(1.2), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    badge.adjustments[0] = 0.3
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = part
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.3), Inches(0.7),
                 title, font_size=20, color=WHITE, bold=True)

    # Description
    add_text_box(slide, x + Inches(0.2), y + Inches(1.8), Inches(3.3), Inches(1.5),
                 desc, font_size=14, color=LIGHT_GRAY)

    # Slides range
    add_text_box(slide, x + Inches(0.2), y + Inches(3.8), Inches(3.3), Inches(0.4),
                 slides_range, font_size=12, color=MED_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 3: PART 1 — SCHISM BMI: THE BIG PICTURE
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)
add_section_number(slide, "1")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "SCHISM's BMI: Inline #ifdef Approach", font_size=30, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(12), Inches(0.5),
             "SCHISM does NOT have a separate BMI wrapper file — BMI is ENABLED directly inside model source code",
             font_size=16, color=ACCENT_YELLOW, bold=True)

# Key point card
add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.1), Inches(0.4),
             "What SCHISM Does", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.4), Inches(5.1), Inches(1.8), [
    "Uses #ifdef USE_NWM_BMI preprocessor flags",
    "Flags are scattered across 3 existing source files",
    "CMake option: define_opt(USE_NWM_BMI ... OFF)",
    "When OFF: SCHISM compiles as standalone model",
    "When ON: Extra BMI-compatible code paths activate",
], font_size=13)

# Files card
add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.5), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.3), Inches(0.4),
             "Files With #ifdef USE_NWM_BMI", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.4), Inches(5.3), Inches(1.8), [
    "schism_init.F90 — initialization checks (1 block)",
    "schism_step.F90 — timestep coupling (2 blocks)",
    "misc_subs.F90 — helper boundary conditions (2 blocks)",
    "Total: Only ~5 #ifdef blocks across 3 files",
    "No separate bmischism.f90 wrapper exists!",
], font_size=13, bullet_color=ACCENT_GREEN)

# Code example
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.4),
             "Example from schism_init.F90 (line 1141):", font_size=12, color=MED_GRAY)
add_code_block(slide, Inches(0.8), Inches(4.9), Inches(5.5), Inches(1.5),
"""! Inside SCHISM's own source code:
#ifdef USE_NWM_BMI
  if(if_source==0) then
    call parallel_abort('USE_NWM_BMI needs if_source>0')
  endif
#endif""", font_size=11)

# ML Analogy card
add_shape(slide, Inches(6.8), Inches(4.5), Inches(5.7), Inches(1.9), CARD_BG, ACCENT_PURPLE, Pt(1))
add_text_box(slide, Inches(7.0), Inches(4.6), Inches(5.3), Inches(0.4),
             "ML Analogy", font_size=16, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(7.0), Inches(5.1), Inches(5.3), Inches(1.2),
             "Like adding @torch.jit.export decorators directly\n"
             "inside your model.py file — the model itself knows\n"
             "how to export, no separate export script needed.\n\n"
             "if USE_ONNX:  # conditional compilation\n"
             "    model.register_buffer('export_ready', ...)",
             font_size=13, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 4: SCHISM ARCHITECTURE — WHY INLINE WORKS
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)
add_section_number(slide, "2")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "Why Inline #ifdef Works for SCHISM", font_size=30, color=WHITE, bold=True)

# SCHISM Architecture - Left side
add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(5.1), Inches(0.4),
             "SCHISM's Simple Architecture", font_size=18, color=ACCENT_BLUE, bold=True)

add_code_block(slide, Inches(1.0), Inches(1.8), Inches(5.1), Inches(3.0),
"""SCHISM Call Stack (Shallow = Easy to #ifdef):

schism_driver.F90 (180 lines)
  |
  +-- schism_init()    <-- #ifdef here (1 block)
  |
  +-- DO timestep = 1, nsteps
  |     |
  |     +-- schism_step()  <-- #ifdef here (2 blocks)
  |
  +-- schism_finalize()

Total depth: 3 levels
Total files modified: 3
Total #ifdef blocks: 5""", font_size=11)

add_bullet_list(slide, Inches(1.0), Inches(5.1), Inches(5.1), Inches(1.3), [
    "Single unstructured triangular grid",
    "One main time loop (simple DO loop)",
    "Already has init/step/finalize separation",
    "Shallow call stack = easy to intercept",
], font_size=13)

# Right side: 4 reasons
add_shape(slide, Inches(6.8), Inches(1.2), Inches(5.7), Inches(5.5), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.3), Inches(0.4),
             "4 Reasons Inline BMI Works Here", font_size=18, color=ACCENT_GREEN, bold=True)

reasons = [
    ("1. Single Grid", "SCHISM uses ONE unstructured mesh.\nBMI grid functions map 1:1 — no branching needed."),
    ("2. Simple Driver", "schism_driver.F90 is only 180 lines.\nAlready structured as init → loop → finalize."),
    ("3. Shallow Nesting", "main → schism_step() is only 2 levels deep.\n#ifdef blocks don't create confusion."),
    ("4. Minimal Coupling", "No multi-model orchestration inside.\nSCHISM is one monolithic physics engine."),
]

for i, (title, desc) in enumerate(reasons):
    y_off = Inches(1.9 + i * 1.25)
    add_text_box(slide, Inches(7.2), y_off, Inches(5.1), Inches(0.35),
                 title, font_size=15, color=ACCENT_YELLOW, bold=True)
    add_text_box(slide, Inches(7.2), y_off + Inches(0.35), Inches(5.1), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 5: SCHISM CMake Configuration
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)
add_section_number(slide, "3")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "SCHISM: How #ifdef Gets Compiled", font_size=30, color=WHITE, bold=True)

# Left: CMake flow
add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(5.1), Inches(0.4),
             "Build Process Flow", font_size=18, color=ACCENT_BLUE, bold=True)

add_code_block(slide, Inches(1.0), Inches(1.8), Inches(5.1), Inches(1.8),
"""# In CMakeLists.txt:
define_opt(USE_NWM_BMI "Use NWM BMI" OFF)

# To enable BMI at build time:
cmake .. -DUSE_NWM_BMI=ON

# This adds: -DUSE_NWM_BMI to compiler flags
# Which activates all #ifdef USE_NWM_BMI blocks""", font_size=11)

add_code_block(slide, Inches(1.0), Inches(3.9), Inches(5.1), Inches(2.5),
"""      COMPILE WITHOUT BMI          COMPILE WITH BMI
      ==================          =================
      cmake ..                    cmake .. -DUSE_NWM_BMI=ON
           |                           |
      #ifdef blocks SKIPPED       #ifdef blocks INCLUDED
           |                           |
      Standard SCHISM             SCHISM + BMI coupling
      (standalone mode)           (can receive external data)

Same source code, different behavior!""", font_size=11)

# Right: What gets activated
add_shape(slide, Inches(6.8), Inches(1.2), Inches(5.7), Inches(5.5), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.3), Inches(0.4),
             "What the #ifdef Blocks Actually Do", font_size=18, color=ACCENT_GREEN, bold=True)

blocks = [
    ("schism_init.F90", "Validation check — ensures if_source > 0\nso external BMI can inject source/sink data"),
    ("schism_step.F90 (block 1)", "Receives discharge data from WRF-Hydro\nvia BMI coupling variables at river mouths"),
    ("schism_step.F90 (block 2)", "Sends back water elevation data to\nWRF-Hydro for coastal feedback"),
    ("misc_subs.F90 (2 blocks)", "Adjusts boundary condition handling\nwhen external BMI data is available"),
]

for i, (file, desc) in enumerate(blocks):
    y_off = Inches(1.9 + i * 1.2)
    add_text_box(slide, Inches(7.2), y_off, Inches(5.1), Inches(0.3),
                 file, font_size=13, color=ACCENT_YELLOW, bold=True)
    add_text_box(slide, Inches(7.2), y_off + Inches(0.3), Inches(5.1), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GRAY)

add_shape(slide, Inches(7.0), Inches(6.0), Inches(5.3), Inches(0.5), RGBColor(0x2C, 0x3E, 0x50))
add_text_box(slide, Inches(7.2), Inches(6.05), Inches(5.0), Inches(0.4),
             "Key Insight: SCHISM's BMI is for COUPLING only, not full 41-function BMI",
             font_size=11, color=ACCENT_ORANGE, bold=True)


# ─────────────────────────────────────────────────
# SLIDE 6: SCHISM SUMMARY
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)
add_section_number(slide, "4")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "SCHISM BMI: Summary & Key Takeaways", font_size=30, color=WHITE, bold=True)

# Visual diagram
add_shape(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(2.5), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(5), Inches(0.4),
             "SCHISM's Approach: Inline Conditional Compilation", font_size=18, color=ACCENT_BLUE, bold=True)

add_code_block(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.6),
"""  [SCHISM Source Code]                            [Build Output]
  ====================                            ==============
  schism_init.F90  ----+
  schism_step.F90  ----|-- cmake -DUSE_NWM_BMI=ON ---->  schism.exe (with BMI coupling)
  misc_subs.F90    ----+

  schism_init.F90  ----+
  schism_step.F90  ----|-- cmake (default OFF)    ---->  schism.exe (standalone)
  misc_subs.F90    ----+""", font_size=11)

# Takeaway cards
takeaways = [
    ("No Separate Wrapper", "SCHISM has NO bmischism.f90 file.\nBMI code lives INSIDE model files.", ACCENT_RED),
    ("Not Full BMI", "Only enables data exchange hooks.\nNot all 41 BMI functions implemented.", ACCENT_ORANGE),
    ("Compile-Time Toggle", "#ifdef = if/else at compile time.\nSame code, two build modes.", ACCENT_BLUE),
    ("Works Because Simple", "Shallow call stack, single grid,\nalready has init/step/finalize.", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(takeaways):
    x = Inches(0.8 + i * 3.05)
    y = Inches(4.1)
    add_shape(slide, x, y, Inches(2.8), Inches(2.2), CARD_BG, color, Pt(1.5))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.5), Inches(0.4),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.65), Inches(2.5), Inches(1.3),
                 desc, font_size=13, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 7: PART 2 INTRO — WRF-HYDRO
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_ORANGE)
add_section_number(slide, "5")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "WRF-Hydro: Can We Use the Same Approach?", font_size=30, color=WHITE, bold=True)

# Big NO
add_shape(slide, Inches(3.5), Inches(1.5), Inches(6.3), Inches(2.0), CARD_BG, ACCENT_RED, Pt(2))
add_text_box(slide, Inches(3.7), Inches(1.6), Inches(5.9), Inches(0.7),
             "Short Answer: NO", font_size=36, color=ACCENT_RED, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.7), Inches(2.3), Inches(5.9), Inches(1.0),
             "WRF-Hydro is fundamentally different from SCHISM in architecture.\n"
             "Inline #ifdef would require modifying 10+ files and create unmaintainable code.",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 5 reasons preview
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(12), Inches(0.4),
             "5 Reasons Why (Detailed in Next Slides):", font_size=18, color=ACCENT_ORANGE, bold=True)

reasons_preview = [
    ("1", "Deep Nesting", "5 levels deep call stack vs SCHISM's 3", ACCENT_RED),
    ("2", "Multi-Component", "5 sequential physics phases per timestep", ACCENT_ORANGE),
    ("3", "3 Grid Types", "1km + 250m + channel network", ACCENT_YELLOW),
    ("4", "235 Source Files", "vs SCHISM's simple structure", ACCENT_PURPLE),
    ("5", "Integrated Time Loop", "No clean init/step/finalize split", ACCENT_BLUE),
]

for i, (num, title, desc, color) in enumerate(reasons_preview):
    x = Inches(0.8 + i * 2.45)
    y = Inches(4.7)
    add_shape(slide, x, y, Inches(2.2), Inches(2.0), CARD_BG, color, Pt(1))
    # Number circle
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.1), y + Inches(0.15),
                                   Inches(0.4), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.6), y + Inches(0.15), Inches(1.5), Inches(0.35),
                 title, font_size=13, color=color, bold=True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.7), Inches(2.0), Inches(1.0),
                 desc, font_size=11, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 8: WRF-HYDRO ARCHITECTURE DEEP DIVE
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_ORANGE)
add_section_number(slide, "6")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "WRF-Hydro: Deep Nested Architecture", font_size=30, color=WHITE, bold=True)

# Call stack visualization
add_shape(slide, Inches(0.8), Inches(1.2), Inches(7.0), Inches(5.8), CARD_BG, ACCENT_ORANGE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(6.6), Inches(0.4),
             "WRF-Hydro Call Stack (5 Levels Deep!)", font_size=18, color=ACCENT_ORANGE, bold=True)

add_code_block(slide, Inches(1.0), Inches(1.8), Inches(6.6), Inches(5.0),
"""main_hrldas_driver.F  (42 lines)       <- Level 1
  |
  +-- land_driver_ini()                   <- INIT
  |
  +-- DO ITIME = 1, NTIME                <- Time Loop
  |     |
  |     +-- land_driver_exe()  (2,869 ln) <- Level 2
  |           |
  |           +-- Read forcing data        Phase 1
  |           +-- Run NoahMP LSM           Phase 2
  |           +-- WTABLE (groundwater)     Phase 3
  |           +-- hrldas_drv_HYDRO()       Phase 4
  |           |     |                      <- Level 3
  |           |     +-- hrldas_cpl_HYDRO()
  |           |           |                <- Level 4
  |           |           +-- HYDRO_exe()
  |           |                 |          <- Level 5
  |           |                 +-- disaggregateDomain
  |           |                 +-- SubsurfaceRouting
  |           |                 +-- OverlandRouting
  |           |                 +-- GW Baseflow
  |           |                 +-- ChannelRouting
  |           |                 +-- aggregateDomain
  |           |
  |           +-- Write output             Phase 5
  |
  +-- land_driver_finish()                <- FINALIZE""", font_size=10)

# Right side: comparison
add_shape(slide, Inches(8.2), Inches(1.2), Inches(4.5), Inches(2.5), CARD_BG, ACCENT_RED, Pt(1))
add_text_box(slide, Inches(8.4), Inches(1.3), Inches(4.1), Inches(0.4),
             "Compare: SCHISM vs WRF-Hydro", font_size=16, color=ACCENT_RED, bold=True)

add_table(slide, Inches(8.4), Inches(1.8), Inches(4.1), 5, 3, [
    ["Metric", "SCHISM", "WRF-Hydro"],
    ["Call depth", "3 levels", "5 levels"],
    ["Driver lines", "180", "2,869"],
    ["Source files", "~50", "235"],
    ["Grids", "1", "3"],
], col_widths=[Inches(1.3), Inches(1.2), Inches(1.6)])

# Problem highlight
add_shape(slide, Inches(8.2), Inches(4.0), Inches(4.5), Inches(3.0), CARD_BG, ACCENT_YELLOW, Pt(1))
add_text_box(slide, Inches(8.4), Inches(4.1), Inches(4.1), Inches(0.4),
             "The Problem:", font_size=16, color=ACCENT_YELLOW, bold=True)
add_text_box(slide, Inches(8.4), Inches(4.6), Inches(4.1), Inches(2.2),
             "HYDRO_exe() is buried 5 levels deep.\n\n"
             "To add #ifdef BMI here, you'd need to\n"
             "modify every level above it to pass\n"
             "BMI control flags downward.\n\n"
             "That means touching 4+ driver files\n"
             "and breaking WRF-Hydro's own code.",
             font_size=13, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 9: WRF-HYDRO 5 PHASES PER TIMESTEP
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_ORANGE)
add_section_number(slide, "7")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "WRF-Hydro: 5 Phases Per Timestep", font_size=30, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.4),
             "Each call to land_driver_exe() runs 5 sequential physics phases with data dependencies:",
             font_size=14, color=LIGHT_GRAY)

phases = [
    ("Phase 1", "Read Forcing", "Rain, temperature, wind,\nhumidity, radiation\nfrom LDASIN files", ACCENT_BLUE, "Input"),
    ("Phase 2", "NoahMP LSM", "Soil moisture, evaporation,\nsnow processes, vegetation\non 1km grid", ACCENT_GREEN, "Physics"),
    ("Phase 3", "Water Table", "Groundwater coupling,\nWTABLE_MMF_NOAHMP()\nsoil-aquifer exchange", ACCENT_PURPLE, "Physics"),
    ("Phase 4", "HYDRO_exe()", "Subsurface + Overland +\nChannel + Baseflow routing\non 250m grid", ACCENT_ORANGE, "Routing"),
    ("Phase 5", "Output", "Write LDASOUT, CHRTOUT,\nRTOUT, restart files\nto NetCDF", ACCENT_RED, "I/O"),
]

for i, (phase, name, desc, color, category) in enumerate(phases):
    x = Inches(0.5 + i * 2.55)
    y = Inches(1.6)

    # Phase card
    add_shape(slide, x, y, Inches(2.3), Inches(3.0), CARD_BG, color, Pt(1.5))

    # Phase number
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(1.0), Inches(0.3),
                 phase, font_size=11, color=MED_GRAY)

    # Category badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x + Inches(1.3), y + Inches(0.1), Inches(0.85), Inches(0.3))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    badge.adjustments[0] = 0.3
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = category
    p.font.size = Pt(9)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Name
    add_text_box(slide, x + Inches(0.1), y + Inches(0.5), Inches(2.1), Inches(0.5),
                 name, font_size=16, color=color, bold=True)

    # Description
    add_text_box(slide, x + Inches(0.1), y + Inches(1.1), Inches(2.1), Inches(1.5),
                 desc, font_size=12, color=LIGHT_GRAY)

    # Arrow between phases
    if i < 4:
        arrow_x = x + Inches(2.35)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       arrow_x, y + Inches(1.2), Inches(0.2), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MED_GRAY
        arrow.line.fill.background()

# Bottom: key insight
add_shape(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8), CARD_BG, ACCENT_RED, Pt(1.5))
add_text_box(slide, Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.4),
             "Why This Makes Inline BMI Impossible:", font_size=16, color=ACCENT_RED, bold=True)
add_text_box(slide, Inches(1.0), Inches(5.6), Inches(5.3), Inches(1.0),
             "BMI's update() must do exactly ONE timestep.\n"
             "But WRF-Hydro's single timestep = 5 phases\n"
             "with data flowing from each phase to the next.\n"
             "Phase 4 DEPENDS on Phase 2 output!",
             font_size=13, color=LIGHT_GRAY)

add_text_box(slide, Inches(6.5), Inches(5.6), Inches(5.5), Inches(1.0),
             "ML Analogy: Like a training step that does\n"
             "data_loading → forward_pass → loss_calc → backward → optimizer\n"
             "You can't just #ifdef one step — they're sequential!\n"
             "You need a wrapper that orchestrates all 5 phases.",
             font_size=13, color=ACCENT_PURPLE)


# ─────────────────────────────────────────────────
# SLIDE 10: WRF-HYDRO 3 GRIDS PROBLEM
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_ORANGE)
add_section_number(slide, "8")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "WRF-Hydro: 3 Simultaneous Grids", font_size=30, color=WHITE, bold=True)

# Grid 0
add_shape(slide, Inches(0.8), Inches(1.2), Inches(3.7), Inches(4.0), CARD_BG, ACCENT_BLUE, Pt(1.5))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(3.3), Inches(0.3),
             "Grid 0: LSM Grid", font_size=18, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(3.3), Inches(0.3),
             "uniform_rectilinear", font_size=13, color=ACCENT_YELLOW)
add_bullet_list(slide, Inches(1.0), Inches(2.1), Inches(3.3), Inches(2.5), [
    "Resolution: 1 km",
    "Variables: soil moisture, snow,\n  temperature, evaporation",
    "NoahMP land surface model",
    "Grid dimensions: ix x jx",
    "BMI type: uniform_rectilinear",
], font_size=12)

# Grid 1
add_shape(slide, Inches(4.8), Inches(1.2), Inches(3.7), Inches(4.0), CARD_BG, ACCENT_GREEN, Pt(1.5))
add_text_box(slide, Inches(5.0), Inches(1.3), Inches(3.3), Inches(0.3),
             "Grid 1: Routing Grid", font_size=18, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(5.0), Inches(1.7), Inches(3.3), Inches(0.3),
             "uniform_rectilinear", font_size=13, color=ACCENT_YELLOW)
add_bullet_list(slide, Inches(5.0), Inches(2.1), Inches(3.3), Inches(2.5), [
    "Resolution: 250 m (4x finer)",
    "Variables: surface head,\n  subsurface flow, overland flow",
    "Terrain routing engine",
    "Grid dimensions: ixrt x jxrt",
    "BMI type: uniform_rectilinear",
], font_size=12, bullet_color=ACCENT_GREEN)

# Grid 2
add_shape(slide, Inches(8.8), Inches(1.2), Inches(3.7), Inches(4.0), CARD_BG, ACCENT_ORANGE, Pt(1.5))
add_text_box(slide, Inches(9.0), Inches(1.3), Inches(3.3), Inches(0.3),
             "Grid 2: Channel Network", font_size=18, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(9.0), Inches(1.7), Inches(3.3), Inches(0.3),
             "vector / unstructured", font_size=13, color=ACCENT_YELLOW)
add_bullet_list(slide, Inches(9.0), Inches(2.1), Inches(3.3), Inches(2.5), [
    "~2.7 million reaches (NWM)",
    "Variables: streamflow (QLINK),\n  channel coordinates",
    "Muskingum-Cunge routing",
    "Network: NLINKS reaches",
    "BMI type: unstructured/vector",
], font_size=12, bullet_color=ACCENT_ORANGE)

# Why this is a problem
add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5), CARD_BG, ACCENT_RED, Pt(1))
add_text_box(slide, Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.4),
             "Why 3 Grids Break Inline BMI:", font_size=16, color=ACCENT_RED, bold=True)
add_text_box(slide, Inches(1.0), Inches(6.1), Inches(5.5), Inches(0.8),
             "Every BMI grid function (get_grid_type, get_grid_size,\n"
             "get_grid_x/y, etc.) must branch on grid ID.\n"
             "Inline #ifdef can't handle this — you need a lookup table.",
             font_size=13, color=LIGHT_GRAY)
add_text_box(slide, Inches(6.8), Inches(6.1), Inches(5.0), Inches(0.8),
             "SCHISM comparison: ONE grid (unstructured triangular).\n"
             "Every grid function returns ONE answer.\n"
             "No branching, no lookup table needed.",
             font_size=13, color=ACCENT_BLUE)


# ─────────────────────────────────────────────────
# SLIDE 11: THE SOLUTION — SEPARATE WRAPPER
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_ORANGE)
add_section_number(slide, "9")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "Our Solution: Separate Non-Invasive Wrapper", font_size=30, color=WHITE, bold=True)

# Architecture diagram
add_shape(slide, Inches(0.8), Inches(1.2), Inches(7.5), Inches(5.8), CARD_BG, ACCENT_GREEN, Pt(1.5))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(7.1), Inches(0.4),
             "bmi_wrf_hydro.f90 — Our Wrapper Architecture", font_size=18, color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(1.0), Inches(1.8), Inches(7.1), Inches(5.0),
"""                    [Python / PyMT / Jupyter]
                              |
                    bmi_initialize("config.yaml")
                    bmi_update()
                    bmi_get_value("streamflow", data)
                    bmi_finalize()
                              |
          +-------------------+-------------------+
          |                                       |
  [bmi_wrf_hydro.f90]                   [bmi_wrf_hydro.f90]
  41 BMI Functions                      Variable Mapping
  ====================                  ==================
  initialize() ---> HYDRO_ini()         "channel_water__volume_
  update()     ---> land_driver_exe()     flow_rate" -> QLINK1
  finalize()   ---> HYDRO_finish()      "soil_water__volume_
  get_value()  ---> select case(name)     fraction" -> SOIL_M
  set_value()  ---> direct assignment
  get_grid_*() ---> grid lookup table
          |                                       |
          +-------------------+-------------------+
                              |
              [WRF-Hydro Source Code — UNTOUCHED]
              ====================================
              module_HYDRO_drv.F90  (HYDRO_ini/exe/finish)
              module_NoahMP_hrldas_driver.F  (land_driver)
              module_RT.F90, module_channel_routing.F90 ...""", font_size=10)

# Right side: benefits
add_shape(slide, Inches(8.7), Inches(1.2), Inches(3.9), Inches(5.8), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(8.9), Inches(1.3), Inches(3.5), Inches(0.4),
             "Why This Is Better", font_size=18, color=ACCENT_BLUE, bold=True)

benefits = [
    ("Non-Invasive", "Zero changes to WRF-Hydro\nsource code. Model stays\nas-is from NCAR."),
    ("Clean Separation", "BMI logic in 1 file (~800 lines)\nvs scattered #ifdef in 10+ files."),
    ("Testable", "Can test BMI wrapper\nindependently with a\nsimple Fortran driver."),
    ("Maintainable", "Update WRF-Hydro version?\nJust rebuild — wrapper\ncalls same subroutines."),
    ("Full BMI", "All 41 functions, 3 grids,\nCSDMS Standard Names,\nready for PyMT."),
]

for i, (title, desc) in enumerate(benefits):
    y_off = Inches(1.85 + i * 1.05)
    add_text_box(slide, Inches(8.9), y_off, Inches(3.5), Inches(0.3),
                 f"✓ {title}", font_size=13, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(9.2), y_off + Inches(0.3), Inches(3.2), Inches(0.6),
                 desc, font_size=11, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 12: COMPARISON TABLE
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_ORANGE)
add_section_number(slide, "10")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "Side-by-Side: Inline #ifdef vs Separate Wrapper", font_size=30, color=WHITE, bold=True)

# Main comparison table
add_table(slide, Inches(0.8), Inches(1.2), Inches(11.7), 9, 3, [
    ["Aspect", "Inline #ifdef (SCHISM Style)", "Separate Wrapper (Our Approach)"],
    ["Files modified", "10+ existing WRF-Hydro files", "1 new file (bmi_wrf_hydro.f90)"],
    ["WRF-Hydro code", "MODIFIED — breaks NCAR updates", "UNTOUCHED — just calls subroutines"],
    ["Complexity", "#ifdef spaghetti across 5,200+ lines", "Clean 41-function module (~800 lines)"],
    ["Grid handling", "3-grid branching in every #ifdef", "Lookup table — 1 handler per grid"],
    ["Testing", "Can't test BMI without full build", "Test wrapper independently"],
    ["Maintainability", "Couples BMI to model internals", "Independent — update either side"],
    ["MPI/Parallel", "State-sharing bugs across #ifdef", "Encapsulated — no MPI conflicts"],
    ["BMI Compliance", "Partial (only coupling hooks)", "Full 41-function implementation"],
], col_widths=[Inches(2.0), Inches(4.85), Inches(4.85)])

# Bottom verdict
add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5), CARD_BG, ACCENT_GREEN, Pt(2))
add_text_box(slide, Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.4),
             "Verdict: Separate Wrapper is the ONLY Viable Approach for WRF-Hydro",
             font_size=20, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.7),
             "Following the exact same pattern as bmi_heat.f90 (our learning template), we'll write bmi_wrf_hydro.f90\n"
             "as a Fortran 2003 module that extends the abstract BMI type and calls WRF-Hydro's existing subroutines.",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────
# SLIDE 13: PART 3 — BMI FORTRAN TEMPLATE INTRO
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)
add_section_number(slide, "11")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "BMI Fortran Template: The Abstract Interface", font_size=30, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.4),
             "bmi-fortran/bmi.f90 — 564 lines that define the BMI contract every model must follow",
             font_size=14, color=ACCENT_YELLOW)

# Left: What it is
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5.1), Inches(0.4),
             "What is bmi.f90?", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5.1), Inches(1.6), [
    "Abstract Fortran type (like Python ABC/interface)",
    "Defines 53 deferred procedures (41 logical functions)",
    "Extra procedures: type overloads for int/float/double",
    "Module name: bmif_2_0 (BMI version 2.0)",
    "Published by CSDMS (Hutton et al. 2020)",
], font_size=13)

# Right: ML analogy
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(2.5), CARD_BG, ACCENT_PURPLE, Pt(1))
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.3), Inches(0.4),
             "ML Analogy: Abstract Base Class", font_size=18, color=ACCENT_PURPLE, bold=True)

add_code_block(slide, Inches(7.0), Inches(2.2), Inches(5.3), Inches(1.6),
"""# Python equivalent of bmi.f90:
from abc import ABC, abstractmethod

class BMI(ABC):
    @abstractmethod
    def initialize(self, config: str) -> int: ...
    @abstractmethod
    def update(self) -> int: ...
    @abstractmethod
    def get_value(self, name, dest) -> int: ...
    # ... 38 more abstract methods""", font_size=10)

# The 6 categories
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(12), Inches(0.4),
             "41 BMI Functions in 6 Categories:", font_size=16, color=WHITE, bold=True)

categories = [
    ("Control (4)", "initialize, update,\nupdate_until, finalize", ACCENT_BLUE),
    ("Info (5)", "component_name,\nvar_names, var_count", ACCENT_GREEN),
    ("Var Info (6)", "type, units, grid,\nitemsize, nbytes, location", ACCENT_ORANGE),
    ("Time (5)", "current, start, end,\nstep, units", ACCENT_PURPLE),
    ("Get/Set (6)", "get_value, set_value,\nptr, at_indices", ACCENT_RED),
    ("Grid (17)", "type, rank, size, shape,\nspacing, x/y/z, topology", ACCENT_YELLOW),
]

for i, (name, funcs, color) in enumerate(categories):
    x = Inches(0.8 + i * 2.08)
    y = Inches(4.8)
    add_shape(slide, x, y, Inches(1.88), Inches(2.0), CARD_BG, color, Pt(1))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(1.68), Inches(0.35),
                 name, font_size=13, color=color, bold=True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.55), Inches(1.68), Inches(1.2),
                 funcs, font_size=11, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 14: BMI KEY CONSTANTS & RULES
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)
add_section_number(slide, "12")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "BMI Spec: Key Constants & Rules", font_size=30, color=WHITE, bold=True)

# Left: Constants
add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(3.0), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(5.1), Inches(0.4),
             "BMI Constants (from bmi.f90)", font_size=18, color=ACCENT_BLUE, bold=True)

add_code_block(slide, Inches(1.0), Inches(1.8), Inches(5.1), Inches(2.1),
"""integer, parameter :: BMI_SUCCESS = 0
integer, parameter :: BMI_FAILURE = 1

integer, parameter :: BMI_MAX_COMPONENT_NAME = 2048
integer, parameter :: BMI_MAX_VAR_NAME       = 2048
integer, parameter :: BMI_MAX_TYPE_NAME      = 2048
integer, parameter :: BMI_MAX_UNITS_NAME     = 2048

! Every BMI function returns integer status:
! 0 = success, 1 = failure (or not implemented)""", font_size=11)

# Right: Rules
add_shape(slide, Inches(6.8), Inches(1.2), Inches(5.7), Inches(3.0), CARD_BG, ACCENT_ORANGE, Pt(1))
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.3), Inches(0.4),
             "5 Critical BMI Rules", font_size=18, color=ACCENT_ORANGE, bold=True)

rules = [
    "Every function returns integer (0=success, 1=failure)",
    "Output via intent(out) parameters, NOT return values",
    "Arrays are ALWAYS flattened to 1D (row-major)",
    "Wrapper is NON-INVASIVE — never modify model code",
    "Memory allocation is the MODEL's job, not BMI's",
]
add_bullet_list(slide, Inches(7.0), Inches(1.8), Inches(5.3), Inches(2.0), rules,
                font_size=13, bullet_color=ACCENT_ORANGE)

# Bottom: The abstract type
add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.7), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.4),
             "The Abstract Type Declaration (Simplified)", font_size=16, color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(1.0), Inches(5.1), Inches(5.5), Inches(1.8),
"""module bmif_2_0
  type, abstract :: bmi
  contains
    ! Control
    procedure(bmif_initialize), deferred :: initialize
    procedure(bmif_update), deferred :: update
    procedure(bmif_finalize), deferred :: finalize
    ! ... 50 more deferred procedures

    ! Type overloads (generic interfaces)
    generic :: get_value => get_value_int, &
               get_value_float, get_value_double
  end type
end module""", font_size=10)

add_code_block(slide, Inches(6.8), Inches(5.1), Inches(5.7), Inches(1.8),
"""! Each procedure has a defined interface:
abstract interface
  function bmif_initialize(self, config_file) result(bmi_status)
    import :: bmi
    class(bmi), intent(out) :: self
    character(len=*), intent(in) :: config_file
    integer :: bmi_status
  end function

  function bmif_update(self) result(bmi_status)
    import :: bmi
    class(bmi), intent(inout) :: self
    integer :: bmi_status
  end function
end interface""", font_size=10)


# ─────────────────────────────────────────────────
# SLIDE 15: HEAT MODEL — THE PHYSICS
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)
add_section_number(slide, "13")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "Heat Model: The Simple Physics Model", font_size=30, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.4),
             "heat/heat.f90 — 158 lines. This is the model that bmi_heat.f90 wraps (like WRF-Hydro is what we'll wrap)",
             font_size=14, color=ACCENT_YELLOW)

# Left: The model type
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(3.0), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5.1), Inches(0.4),
             "heat_model Type (= Model State)", font_size=18, color=ACCENT_BLUE, bold=True)

add_code_block(slide, Inches(1.0), Inches(2.2), Inches(5.1), Inches(2.1),
"""type :: heat_model
  integer :: n_x = 0         ! Grid columns
  integer :: n_y = 0         ! Grid rows
  real :: dt = 0.            ! Timestep (seconds)
  real :: t  = 0.            ! Current time
  real :: t_end = 0.         ! End time
  real :: alpha = 0.         ! Diffusivity (physics param)

  ! The main state variable:
  real, pointer :: temperature(:,:) => null()
  ! Temp workspace:
  real, pointer :: temperature_tmp(:,:) => null()
end type""", font_size=11)

# Right: ML analogy
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(3.0), CARD_BG, ACCENT_PURPLE, Pt(1))
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.3), Inches(0.4),
             "ML Analogy: PyTorch Model", font_size=18, color=ACCENT_PURPLE, bold=True)

add_code_block(slide, Inches(7.0), Inches(2.2), Inches(5.3), Inches(2.1),
"""class HeatModel(nn.Module):
    def __init__(self):
        self.n_x = 0          # = Grid width
        self.n_y = 0          # = Grid height
        self.dt = 0.0         # = Learning rate
        self.t = 0.0          # = Current epoch
        self.t_end = 0.0      # = Total epochs
        self.alpha = 0.0      # = Weight decay

        # State = model weights/activations:
        self.temperature = None      # = weights
        self.temperature_tmp = None  # = gradients""", font_size=11)

# Bottom: The 3 key subroutines
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(12), Inches(0.4),
             "3 Model Subroutines (= the entire model API):", font_size=16, color=WHITE, bold=True)

subs = [
    ("initialize_from_file()", "Reads config (4 values),\nallocates temperature(:,:),\nsets initial condition\n(top row = hot boundary)",
     ACCENT_BLUE, "Like model = load_config(yaml)\ntorch.randn(n_y, n_x)"),
    ("advance_in_time()", "2D stencil computation:\nnew = avg(4 neighbors) * alpha\nCopies result, increments t",
     ACCENT_GREEN, "Like model.forward(x):\nconv2d with 3x3 kernel\noptimizer.step()"),
    ("cleanup()", "Deallocates temperature\narrays. Frees memory.",
     ACCENT_RED, "Like del model\ntorch.cuda.empty_cache()"),
]

for i, (name, desc, color, ml) in enumerate(subs):
    x = Inches(0.8 + i * 4.2)
    y = Inches(5.3)
    add_shape(slide, x, y, Inches(3.9), Inches(1.8), CARD_BG, color, Pt(1))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.05), Inches(3.7), Inches(0.3),
                 name, font_size=14, color=color, bold=True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.4), Inches(1.8), Inches(1.2),
                 desc, font_size=10, color=LIGHT_GRAY)
    add_text_box(slide, x + Inches(2.0), y + Inches(0.4), Inches(1.8), Inches(1.2),
                 ml, font_size=10, color=ACCENT_PURPLE)


# ─────────────────────────────────────────────────
# SLIDE 16: BMI HEAT WRAPPER — TYPE EXTENSION
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)
add_section_number(slide, "14")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "bmi_heat.f90: The Wrapper Pattern", font_size=30, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.4),
             "935 lines — THIS is the template for our bmi_wrf_hydro.f90. Every pattern here gets reused.",
             font_size=14, color=ACCENT_YELLOW)

# Type extension
add_shape(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(3.0), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5.6), Inches(0.4),
             "Pattern 1: Type Extension (Inheritance)", font_size=18, color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(1.0), Inches(2.2), Inches(5.6), Inches(2.1),
"""module bmiheatf
  use heatf           ! Import the model
  use bmif_2_0        ! Import BMI spec
  use, intrinsic :: iso_c_binding  ! For pointers

  type, extends(bmi) :: bmi_heat    ! Inherit from bmi
    private
    type(heat_model) :: model       ! EMBED the model
  contains
    procedure :: initialize => heat_initialize
    procedure :: update => heat_update
    procedure :: finalize => heat_finalize
    ! ... all 53 procedure mappings
  end type bmi_heat
end module""", font_size=11)

# ML equivalent
add_shape(slide, Inches(7.2), Inches(1.6), Inches(5.3), Inches(3.0), CARD_BG, ACCENT_PURPLE, Pt(1))
add_text_box(slide, Inches(7.4), Inches(1.7), Inches(4.9), Inches(0.4),
             "ML Equivalent: Wrapper Class", font_size=18, color=ACCENT_PURPLE, bold=True)

add_code_block(slide, Inches(7.4), Inches(2.2), Inches(4.9), Inches(2.1),
"""class BMIHeat(BMI):  # Inherit abstract class
    def __init__(self):
        self._model = HeatModel()  # EMBED model

    def initialize(self, config):
        self._model.load_config(config)
        return BMI_SUCCESS

    def update(self):
        self._model.advance_in_time()
        return BMI_SUCCESS

    def get_value(self, name, dest):
        if name == "temperature":
            dest[:] = self._model.temperature.flatten()
        return BMI_SUCCESS""", font_size=10)

# Key insight
add_shape(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.3), CARD_BG, ACCENT_ORANGE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(4.9), Inches(11.3), Inches(0.4),
             "3 Key Patterns We'll Reuse for WRF-Hydro:", font_size=16, color=ACCENT_ORANGE, bold=True)

patterns = [
    ("Type Extension", "type, extends(bmi) :: bmi_wrf_hydro\nEmbeds WRF-Hydro state as private member", ACCENT_BLUE),
    ("Procedure Mapping", "procedure :: initialize => wrfhydro_initialize\nMaps BMI names to our implementation", ACCENT_GREEN),
    ("Private Model Instance", "type(rt_domain_type) :: model\nWrapper owns the model, callers see only BMI", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(patterns):
    x = Inches(1.0 + i * 3.85)
    y = Inches(5.4)
    add_text_box(slide, x, y, Inches(3.5), Inches(0.3),
                 f"  {title}", font_size=14, color=color, bold=True)
    add_text_box(slide, x, y + Inches(0.35), Inches(3.5), Inches(1.0),
                 desc, font_size=11, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 17: BMI HEAT — CONTROL FUNCTIONS
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)
add_section_number(slide, "15")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "bmi_heat.f90: Control Functions (IRF)", font_size=30, color=WHITE, bold=True)

# Initialize
add_shape(slide, Inches(0.8), Inches(1.2), Inches(3.8), Inches(5.8), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(3.4), Inches(0.4),
             "initialize(config_file)", font_size=16, color=ACCENT_BLUE, bold=True)

add_code_block(slide, Inches(1.0), Inches(1.8), Inches(3.4), Inches(2.5),
"""function heat_initialize(self, &
    config_file) result(bmi_status)
  class(bmi_heat), intent(out) :: self
  character(len=*), intent(in) :: &
      config_file

  ! Just delegate to model:
  call initialize_from_file( &
      self%model, config_file)

  bmi_status = BMI_SUCCESS
end function""", font_size=10)

add_text_box(slide, Inches(1.0), Inches(4.5), Inches(3.4), Inches(2.3),
             "What it does:\n"
             "1. Takes config file path\n"
             "2. Calls model's own init\n"
             "3. Returns success/failure\n\n"
             "For WRF-Hydro:\n"
             "Will call HYDRO_ini() and\n"
             "land_driver_ini()",
             font_size=11, color=LIGHT_GRAY)

# Update
add_shape(slide, Inches(4.9), Inches(1.2), Inches(3.8), Inches(5.8), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(5.1), Inches(1.3), Inches(3.4), Inches(0.4),
             "update()", font_size=16, color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(5.1), Inches(1.8), Inches(3.4), Inches(2.5),
"""function heat_update(self) &
    result(bmi_status)
  class(bmi_heat), intent(inout) :: self

  ! ONE timestep only:
  call advance_in_time(self%model)

  bmi_status = BMI_SUCCESS
end function

! Also: update_until(time)
! Loops update() until target
! time is reached.""", font_size=10)

add_text_box(slide, Inches(5.1), Inches(4.5), Inches(3.4), Inches(2.3),
             "What it does:\n"
             "1. Advances model by dt\n"
             "2. Exactly ONE step\n"
             "3. Caller controls the loop!\n\n"
             "For WRF-Hydro:\n"
             "Will call land_driver_exe()\n"
             "for exactly one timestep",
             font_size=11, color=LIGHT_GRAY)

# Finalize
add_shape(slide, Inches(9.0), Inches(1.2), Inches(3.8), Inches(5.8), CARD_BG, ACCENT_RED, Pt(1))
add_text_box(slide, Inches(9.2), Inches(1.3), Inches(3.4), Inches(0.4),
             "finalize()", font_size=16, color=ACCENT_RED, bold=True)

add_code_block(slide, Inches(9.2), Inches(1.8), Inches(3.4), Inches(2.5),
"""function heat_finalize(self) &
    result(bmi_status)
  class(bmi_heat), intent(inout) :: self

  ! Cleanup:
  call cleanup(self%model)

  bmi_status = BMI_SUCCESS
end function

! cleanup() deallocates
! temperature arrays and
! frees all memory.""", font_size=10)

add_text_box(slide, Inches(9.2), Inches(4.5), Inches(3.4), Inches(2.3),
             "What it does:\n"
             "1. Frees all memory\n"
             "2. Closes files\n"
             "3. Clean shutdown\n\n"
             "For WRF-Hydro:\n"
             "Will call HYDRO_finish()\n"
             "and land_driver_finish()",
             font_size=11, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 18: BMI HEAT — SELECT CASE PATTERN
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)
add_section_number(slide, "16")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "bmi_heat.f90: The select case Pattern", font_size=30, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.4),
             "Every variable-info function uses the same pattern: switch on variable name string",
             font_size=14, color=ACCENT_YELLOW)

# The pattern
add_shape(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.4), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5.6), Inches(0.4),
             "Pattern: String-Based Variable Dispatch", font_size=18, color=ACCENT_BLUE, bold=True)

add_code_block(slide, Inches(1.0), Inches(2.2), Inches(5.6), Inches(4.5),
"""! get_var_type — what datatype is this variable?
function heat_var_type(self, name, type) result(s)
  class(bmi_heat), intent(in) :: self
  character(len=*), intent(in) :: name
  character(len=BMI_MAX_TYPE_NAME), intent(out) :: type
  integer :: s

  select case(name)
  case("plate_surface__temperature")
    type = "real"
    s = BMI_SUCCESS
  case("plate_surface__thermal_diffusivity")
    type = "real"
    s = BMI_SUCCESS
  case("model__identification_number")
    type = "integer"
    s = BMI_SUCCESS
  case default
    type = "-"
    s = BMI_FAILURE    ! Unknown variable!
  end select
end function

! Same pattern for: get_var_units, get_var_grid,
! get_var_itemsize, get_var_nbytes, get_var_location""", font_size=10)

# WRF-Hydro equivalent
add_shape(slide, Inches(7.2), Inches(1.6), Inches(5.3), Inches(5.4), CARD_BG, ACCENT_ORANGE, Pt(1))
add_text_box(slide, Inches(7.4), Inches(1.7), Inches(4.9), Inches(0.4),
             "For Our WRF-Hydro Wrapper:", font_size=18, color=ACCENT_ORANGE, bold=True)

add_code_block(slide, Inches(7.4), Inches(2.2), Inches(4.9), Inches(3.0),
"""! In bmi_wrf_hydro.f90:
select case(name)
case("channel_water__volume_flow_rate")
  type = "double precision"   ! QLINK

case("land_surface_water__depth")
  type = "real"               ! sfcheadrt

case("soil_water__volume_fraction")
  type = "real"               ! SOIL_M

case("snowpack__liquid-equivalent_depth")
  type = "real"               ! SNEQV

case default
  s = BMI_FAILURE
end select""", font_size=10)

add_text_box(slide, Inches(7.4), Inches(5.4), Inches(4.9), Inches(1.3),
             "Key Insight:\n"
             "CSDMS Standard Names are the public API.\n"
             "Internal WRF-Hydro names (QLINK, SOIL_M)\n"
             "are implementation details hidden inside\n"
             "the select case mapping.\n\n"
             "ML Analogy: Like a REST API mapping\n"
             "/api/predictions → model.forward(x)",
             font_size=12, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 19: BMI HEAT — GET/SET VALUE
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)
add_section_number(slide, "17")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "bmi_heat.f90: get_value & set_value", font_size=30, color=WHITE, bold=True)

# get_value
add_shape(slide, Inches(0.8), Inches(1.2), Inches(6.0), Inches(3.0), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(5.6), Inches(0.4),
             "get_value: Copy Data Out (Safe)", font_size=18, color=ACCENT_BLUE, bold=True)

add_code_block(slide, Inches(1.0), Inches(1.8), Inches(5.6), Inches(2.1),
"""function heat_get_double(self, name, dest) result(s)
  class(bmi_heat), intent(in) :: self
  character(len=*), intent(in) :: name
  double precision, intent(inout) :: dest(:)  ! 1D!

  select case(name)
  case("plate_surface__temperature")
    ! FLATTEN 2D -> 1D with reshape:
    dest = reshape(self%model%temperature, [nx*ny])
  case default
    s = BMI_FAILURE
  end select
end function""", font_size=10)

# set_value
add_shape(slide, Inches(7.2), Inches(1.2), Inches(5.3), Inches(3.0), CARD_BG, ACCENT_RED, Pt(1))
add_text_box(slide, Inches(7.4), Inches(1.3), Inches(4.9), Inches(0.4),
             "set_value: Push Data In (Safe)", font_size=18, color=ACCENT_RED, bold=True)

add_code_block(slide, Inches(7.4), Inches(1.8), Inches(4.9), Inches(2.1),
"""function heat_set_double(self, name, src) result(s)
  class(bmi_heat), intent(inout) :: self
  character(len=*), intent(in) :: name
  double precision, intent(in) :: src(:)  ! 1D!

  select case(name)
  case("plate_surface__temperature")
    ! UNFLATTEN 1D -> 2D with reshape:
    self%model%temperature = &
        reshape(src, [ny, nx])
  case default
    s = BMI_FAILURE
  end select
end function""", font_size=10)

# get_value_ptr
add_shape(slide, Inches(0.8), Inches(4.5), Inches(6.0), Inches(2.7), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(1.0), Inches(4.6), Inches(5.6), Inches(0.4),
             "get_value_ptr: Zero-Copy Pointer (Fast)", font_size=18, color=ACCENT_GREEN, bold=True)

add_code_block(slide, Inches(1.0), Inches(5.1), Inches(5.6), Inches(1.8),
"""function heat_get_ptr(self, name, dest_ptr) result(s)
  use iso_c_binding, only: c_ptr, c_loc, c_f_pointer
  type(c_ptr), intent(inout) :: dest_ptr

  select case(name)
  case("plate_surface__temperature")
    ! Return raw memory pointer (no copy!):
    dest_ptr = c_loc(self%model%temperature(1,1))
  end select
  ! Caller uses c_f_pointer to access data
end function""", font_size=10)

# Coupling diagram
add_shape(slide, Inches(7.2), Inches(4.5), Inches(5.3), Inches(2.7), CARD_BG, ACCENT_ORANGE, Pt(1))
add_text_box(slide, Inches(7.4), Inches(4.6), Inches(4.9), Inches(0.4),
             "This Is How Models Couple!", font_size=18, color=ACCENT_ORANGE, bold=True)

add_code_block(slide, Inches(7.4), Inches(5.1), Inches(4.9), Inches(1.8),
"""! PyMT coupling loop (pseudocode):
do while (time < end_time)
  call wrfhydro%update()          ! Run WRF-Hydro
  call wrfhydro%get_value(         ! Get discharge
       "channel_water__volume_flow_rate", Q)
  call schism%set_value(           ! Feed to SCHISM
       "channel_water__volume_flow_rate", Q)
  call schism%update()             ! Run SCHISM
end do""", font_size=10)


# ─────────────────────────────────────────────────
# SLIDE 20: BMI HEAT — GRID FUNCTIONS
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)
add_section_number(slide, "18")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "bmi_heat.f90: Grid Functions (17 of them!)", font_size=30, color=WHITE, bold=True)

# Grid metadata
add_shape(slide, Inches(0.8), Inches(1.2), Inches(6.0), Inches(3.2), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(5.6), Inches(0.4),
             "Grid Metadata Functions", font_size=18, color=ACCENT_BLUE, bold=True)

add_code_block(slide, Inches(1.0), Inches(1.8), Inches(5.6), Inches(2.3),
"""! Heat model has 1 grid (id=0): uniform_rectilinear

get_grid_type(0)    => "uniform_rectilinear"
get_grid_rank(0)    => 2         (2D grid)
get_grid_size(0)    => n_x * n_y (total cells)
get_grid_shape(0)   => [n_y, n_x]  (rows, cols)
get_grid_spacing(0) => [1.0, 1.0]  (dx, dy)
get_grid_origin(0)  => [0.0, 0.0]  (x0, y0)

! Topology functions return BMI_FAILURE:
get_grid_edge_count(0)  => BMI_FAILURE (N/A)
get_grid_face_count(0)  => BMI_FAILURE (N/A)
! (These are for unstructured grids only)""", font_size=10)

# WRF-Hydro grids
add_shape(slide, Inches(7.2), Inches(1.2), Inches(5.3), Inches(3.2), CARD_BG, ACCENT_ORANGE, Pt(1))
add_text_box(slide, Inches(7.4), Inches(1.3), Inches(4.9), Inches(0.4),
             "For WRF-Hydro: 3 Grids!", font_size=18, color=ACCENT_ORANGE, bold=True)

add_code_block(slide, Inches(7.4), Inches(1.8), Inches(4.9), Inches(2.3),
"""! Grid 0: LSM (1km)
get_grid_type(0) => "uniform_rectilinear"
get_grid_shape(0) => [jx, ix]
get_grid_spacing(0) => [1000.0, 1000.0]

! Grid 1: Routing (250m)
get_grid_type(1) => "uniform_rectilinear"
get_grid_shape(1) => [jxrt, ixrt]
get_grid_spacing(1) => [250.0, 250.0]

! Grid 2: Channel network
get_grid_type(2) => "unstructured"
get_grid_size(2) => NLINKS  (num reaches)
get_grid_x/y(2) => CHLON(:), CHLAT(:)""", font_size=10)

# Bottom: What each grid function does
add_shape(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.4), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.4),
             "Grid Function Quick Reference:", font_size=16, color=ACCENT_GREEN, bold=True)

add_table(slide, Inches(1.0), Inches(5.3), Inches(11.3), 5, 4, [
    ["Function", "Returns", "Rectilinear Grid", "Unstructured Grid"],
    ["get_grid_type", "string", "\"uniform_rectilinear\"", "\"unstructured\""],
    ["get_grid_rank", "integer", "2 (rows, cols)", "1 (node list)"],
    ["get_grid_shape", "int array", "[n_rows, n_cols]", "BMI_FAILURE"],
    ["get_grid_x/y", "real array", "BMI_FAILURE (use spacing)", "Lat/Lon coordinates"],
], col_widths=[Inches(2.2), Inches(1.6), Inches(3.75), Inches(3.75)])


# ─────────────────────────────────────────────────
# SLIDE 21: THE DRIVER PROGRAM
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)
add_section_number(slide, "19")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "bmi_main.f90: The BMI Driver (Test Program)", font_size=30, color=WHITE, bold=True)

# Full driver code
add_shape(slide, Inches(0.8), Inches(1.2), Inches(6.0), Inches(5.8), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(5.6), Inches(0.4),
             "Complete Driver: 65 Lines", font_size=18, color=ACCENT_BLUE, bold=True)

add_code_block(slide, Inches(1.0), Inches(1.8), Inches(5.6), Inches(5.0),
"""program bmi_main
  use bmiheatf          ! Import BMI wrapper
  type(bmi_heat) :: m   ! Create instance

  ! 1. INITIALIZE
  s = m%initialize(config_file)

  ! 2. QUERY METADATA
  s = m%get_component_name(name)
  s = m%get_input_var_names(names)
  s = m%get_output_var_names(names)

  s = m%get_start_time(time)
  s = m%get_end_time(end_time)
  s = m%get_time_step(dt)

  allocate(temperature(grid_size))

  ! 3. TIME LOOP (caller controls!)
  do while (current_time < end_time)
    s = m%get_value("plate_surface__"//&
        "temperature", temperature)
    print *, temperature(:5)   ! first 5 values

    s = m%update()  ! ONE timestep

    s = m%get_current_time(current_time)
  end do

  ! 4. FINALIZE
  s = m%finalize()
end program""", font_size=10)

# Right: Annotations
add_shape(slide, Inches(7.2), Inches(1.2), Inches(5.3), Inches(5.8), CARD_BG, ACCENT_PURPLE, Pt(1))
add_text_box(slide, Inches(7.4), Inches(1.3), Inches(4.9), Inches(0.4),
             "Key Observations", font_size=18, color=ACCENT_PURPLE, bold=True)

observations = [
    ("The caller controls the time loop",
     "Not the model! BMI inverts control.\n"
     "model.update() = one step only."),
    ("Config file = single argument",
     "initialize('config.cfg') reads all\n"
     "settings. For WRF-Hydro: points to\n"
     "namelist.hrldas + hydro.namelist."),
    ("Only BMI API calls used",
     "Driver never accesses model internals.\n"
     "Everything through get_value/set_value.\n"
     "This is the ENCAPSULATION principle."),
    ("Status checks on every call",
     "Every BMI function returns int status.\n"
     "s = BMI_SUCCESS (0) or BMI_FAILURE (1).\n"
     "Production code should check these!"),
    ("ML Analogy: Training Script",
     "This is like your train.py that calls\n"
     "model.forward(), loss.backward(),\n"
     "optimizer.step() in a loop."),
]

for i, (title, desc) in enumerate(observations):
    y_off = Inches(1.9 + i * 1.05)
    add_text_box(slide, Inches(7.4), y_off, Inches(4.9), Inches(0.3),
                 title, font_size=13, color=ACCENT_YELLOW, bold=True)
    add_text_box(slide, Inches(7.6), y_off + Inches(0.3), Inches(4.7), Inches(0.65),
                 desc, font_size=10, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────
# SLIDE 22: TEST SUITE OVERVIEW
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_GREEN)
add_section_number(slide, "20")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "bmi-example-fortran: Test Suite (52 Tests)", font_size=30, color=WHITE, bold=True)

# Test categories
add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(3.2), CARD_BG, ACCENT_BLUE, Pt(1))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(5.1), Inches(0.4),
             "42 Unit Tests (All Pass!)", font_size=18, color=ACCENT_BLUE, bold=True)

add_table(slide, Inches(1.0), Inches(1.8), Inches(5.1), 8, 2, [
    ["Test File", "What It Tests"],
    ["test_initialize", "Init from config, verify grid dims"],
    ["test_get_component_name", "Component name string"],
    ["test_get_value", "Copy values to 1D array"],
    ["test_set_value", "Modify model state via BMI"],
    ["test_get_value_ptr", "Zero-copy pointer access"],
    ["test_by_reference", "Pointer reflects model changes"],
    ["test_update", "Single timestep + time advance"],
], col_widths=[Inches(2.3), Inches(2.8)])

# Examples
add_shape(slide, Inches(6.8), Inches(1.2), Inches(5.7), Inches(3.2), CARD_BG, ACCENT_GREEN, Pt(1))
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.3), Inches(0.4),
             "7 Example Programs", font_size=18, color=ACCENT_GREEN, bold=True)

add_table(slide, Inches(7.0), Inches(1.8), Inches(5.3), 8, 2, [
    ["Example", "Demonstrates"],
    ["info_ex", "Query component name & var names"],
    ["irf_ex", "Full init -> run -> finalize lifecycle"],
    ["vargrid_ex", "All grid/variable metadata"],
    ["get_value_ex", "get_value, ptr, at_indices"],
    ["set_value_ex", "set_value, at_indices"],
    ["conflicting_instances", "Two models don't interfere"],
    ["change_diffusivity", "Modify params between runs"],
], col_widths=[Inches(2.3), Inches(3.0)])

# Build results
add_shape(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.3), CARD_BG, ACCENT_GREEN, Pt(1.5))
add_text_box(slide, Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.4),
             "Our Build & Test Results: ALL 52 PASSED", font_size=20, color=ACCENT_GREEN, bold=True,
             alignment=PP_ALIGN.CENTER)

add_code_block(slide, Inches(1.0), Inches(5.3), Inches(5.3), Inches(1.4),
"""$ bash build_and_test.sh

Step 1: Conda activate wrfhydro-bmi .. OK
Step 2: CMake build ................ OK
Step 3: Standalone heat model ...... OK
Step 4: BMI driver (bmi_main) ..... OK
Step 5: CTest (49/49) ............. OK
Step 6: Unit tests (42/42) ........ OK
Step 7: Examples (7/7) ............ OK""", font_size=11)

add_code_block(slide, Inches(6.8), Inches(5.3), Inches(5.5), Inches(1.4),
"""FINAL RESULTS:
  Standalone model:     PASS
  BMI driver:           PASS
  CTest suite:       49/49 PASS
  Individual tests:  42/42 PASS
  Example programs:   7/7  PASS

  TOTAL: 52/52 ALL TESTS PASSED""", font_size=11)


# ─────────────────────────────────────────────────
# SLIDE 23: WHAT'S NEXT — ROADMAP
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_YELLOW)
add_section_number(slide, "21")

add_text_box(slide, Inches(1.1), Inches(0.3), Inches(10), Inches(0.6),
             "What's Next: Our Implementation Roadmap", font_size=30, color=WHITE, bold=True)

# Completed items
add_shape(slide, Inches(0.8), Inches(1.2), Inches(3.8), Inches(5.5), CARD_BG, ACCENT_GREEN, Pt(1.5))
add_text_box(slide, Inches(1.0), Inches(1.3), Inches(3.4), Inches(0.4),
             "DONE (Phase 1a)", font_size=18, color=ACCENT_GREEN, bold=True)

done_items = [
    "Study BMI spec (bmi.f90)",
    "Study Heat Model example",
    "Build & test bmi-example-fortran",
    "Study SCHISM BMI approach",
    "Compile WRF-Hydro v5.4.0",
    "Run Croton NY test case",
    "Study WRF-Hydro internals",
    "Create 8 documentation files",
]
add_bullet_list(slide, Inches(1.0), Inches(1.8), Inches(3.4), Inches(4.5), done_items,
                font_size=12, bullet_color=ACCENT_GREEN)

# Next up
add_shape(slide, Inches(4.9), Inches(1.2), Inches(3.8), Inches(5.5), CARD_BG, ACCENT_ORANGE, Pt(1.5))
add_text_box(slide, Inches(5.1), Inches(1.3), Inches(3.4), Inches(0.4),
             "NEXT (Phase 1b)", font_size=18, color=ACCENT_ORANGE, bold=True)

next_items = [
    "Identify 5-10 WRF-Hydro variables",
    "Map to CSDMS Standard Names",
    "Refactor time loop (IRF pattern)",
    "Write bmi_wrf_hydro.f90",
    "Implement 41 BMI functions",
    "Compile as libwrfhydro_bmi.so",
    "Write Fortran test driver",
    "Validate against standalone run",
]
add_bullet_list(slide, Inches(5.1), Inches(1.8), Inches(3.4), Inches(4.5), next_items,
                font_size=12, bullet_color=ACCENT_ORANGE)

# Future
add_shape(slide, Inches(9.0), Inches(1.2), Inches(3.8), Inches(5.5), CARD_BG, ACCENT_PURPLE, Pt(1.5))
add_text_box(slide, Inches(9.2), Inches(1.3), Inches(3.4), Inches(0.4),
             "FUTURE (Phase 2-4)", font_size=18, color=ACCENT_PURPLE, bold=True)

future_items = [
    "Babelize WRF-Hydro (pymt_wrfhydro)",
    "Babelize SCHISM (pymt_schism)",
    "Register PyMT plugins",
    "Run bmi-tester on both",
    "Write coupling script",
    "Grid mapping (1km -> triangles)",
    "Time synchronization",
    "Compound flooding case study",
]
add_bullet_list(slide, Inches(9.2), Inches(1.8), Inches(3.4), Inches(4.5), future_items,
                font_size=12, bullet_color=ACCENT_PURPLE)


# ─────────────────────────────────────────────────
# SLIDE 24: THANK YOU / SUMMARY
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.0),
             "Key Takeaways", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Three key points
takeaway_items = [
    ("SCHISM: Inline #ifdef",
     "Works because SCHISM is simple:\nsingle grid, shallow nesting,\nalready has init/step/finalize.\nOnly ~5 #ifdef blocks in 3 files.",
     ACCENT_BLUE),
    ("WRF-Hydro: Separate Wrapper",
     "Required because WRF-Hydro is complex:\n3 grids, 5-level nesting, 235 files,\n5 phases per timestep.\nWe'll write bmi_wrf_hydro.f90 (~800 lines).",
     ACCENT_ORANGE),
    ("BMI Heat = Our Blueprint",
     "935-line template shows every pattern:\ntype extension, select case dispatch,\nreshape flattening, zero-copy pointers.\n52/52 tests pass in our environment.",
     ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(takeaway_items):
    x = Inches(0.8 + i * 4.2)
    y = Inches(2.5)
    add_shape(slide, x, y, Inches(3.9), Inches(3.5), CARD_BG, color, Pt(2))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.5), Inches(2.3),
                 desc, font_size=14, color=LIGHT_GRAY)

# Bottom
add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), ACCENT_BLUE)
add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
             "WRF-Hydro BMI Project — Phase 1 Progress Report",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════
output_path = "/mnt/c/Users/mohansai/Desktop/Projects/VS_Code/WRF-Hydro-BMI/bmi_wrf_hydro/Docs/BMI_Architecture_Strategy.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
