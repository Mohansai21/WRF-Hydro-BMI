#!/usr/bin/env python3
"""
Generate Weekly Reporting PPT — February 20, 2026
Light theme, Times New Roman, for Professor in Hydrology
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette (Light Theme) ───
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT_BG = RGBColor(0xF0, 0xF2, 0xF5)       # Slide background
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)        # Main text
MED_TEXT = RGBColor(0x33, 0x33, 0x33)         # Body text
LIGHT_TEXT = RGBColor(0x55, 0x55, 0x55)       # Subtle text
BLUE_ACCENT = RGBColor(0x1A, 0x56, 0xDB)     # Primary accent
BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)      # Light blue bg
BLUE_MED = RGBColor(0xD0, 0xE1, 0xFD)        # Medium blue bg
GREEN_ACCENT = RGBColor(0x0D, 0x7C, 0x3E)    # Success green
GREEN_LIGHT = RGBColor(0xE6, 0xF4, 0xEA)     # Light green bg
ORANGE_ACCENT = RGBColor(0xC6, 0x5D, 0x07)   # Warning orange
ORANGE_LIGHT = RGBColor(0xFE, 0xF3, 0xE2)    # Light orange bg
RED_ACCENT = RGBColor(0xC5, 0x22, 0x1F)      # Alert red
RED_LIGHT = RGBColor(0xFC, 0xE8, 0xE6)       # Light red bg
PURPLE_ACCENT = RGBColor(0x6A, 0x1B, 0x9A)   # Purple
PURPLE_LIGHT = RGBColor(0xF3, 0xE5, 0xF5)    # Light purple bg
TEAL_ACCENT = RGBColor(0x00, 0x69, 0x6B)     # Teal
TEAL_LIGHT = RGBColor(0xE0, 0xF2, 0xF1)      # Light teal bg
GRAY_BORDER = RGBColor(0xDA, 0xDC, 0xE0)     # Borders
DARK_CARD = RGBColor(0x26, 0x32, 0x38)        # For code blocks
CODE_GREEN = RGBColor(0x4C, 0xAF, 0x50)       # Code text
TITLE_BG = RGBColor(0x1A, 0x1A, 0x2E)         # Title slide dark bg
TITLE_ACCENT = RGBColor(0x3D, 0x8B, 0xFD)     # Title accent blue

FONT = "Times New Roman"
CODE_FONT = "Consolas"

# Slide dimensions matching template
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0.75)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.03
    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0.75)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=14, color=DARK_TEXT,
             bold=False, alignment=PP_ALIGN.LEFT, font=FONT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, text, size=12, color=MED_TEXT,
                  bold=False, font=FONT, line_spacing=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.space_after = Pt(2)
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def add_bullets(slide, left, top, width, height, items, size=12, color=MED_TEXT,
                bullet_char="\u2022", bullet_color=BLUE_ACCENT, spacing=Pt(4)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.space_before = Pt(1)
        # Bullet
        run_b = p.add_run()
        run_b.text = f"{bullet_char} "
        run_b.font.size = Pt(size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = FONT
        # Text
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(size)
        run_t.font.color.rgb = color
        run_t.font.name = FONT
    return txBox


def add_code(slide, left, top, width, height, code_text, size=9):
    shape = add_shape(slide, left, top, width, height, DARK_CARD, GRAY_BORDER, Pt(0.5))
    txBox = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.08),
                                     width - Inches(0.24), height - Inches(0.16))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = CODE_GREEN
        p.font.name = CODE_FONT
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    return shape


def slide_header(slide, title, subtitle=None):
    """Add standard slide header bar."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.72), BLUE_ACCENT)
    add_text(slide, Inches(0.5), Inches(0.12), Inches(9), Inches(0.5),
             title, size=22, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.72), Inches(9), Inches(0.35),
                 subtitle, size=11, color=LIGHT_TEXT, italic=True)


def add_table(slide, left, top, width, rows, cols, data, col_widths=None,
              header_bg=BLUE_ACCENT, header_fg=WHITE, row_bg1=WHITE, row_bg2=LIGHT_BG):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.32 * rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = header_fg
                else:
                    p.font.color.rgb = MED_TEXT
            cf = cell.fill
            cf.solid()
            if r == 0:
                cf.fore_color.rgb = header_bg
            elif r % 2 == 0:
                cf.fore_color.rgb = row_bg2
            else:
                cf.fore_color.rgb = row_bg1
    return table_shape


# ═══════════════════════════════════════
#  BUILD PRESENTATION
# ═══════════════════════════════════════

prs = Presentation()
prs.slide_width = Emu(9144000)   # 10 inches
prs.slide_height = Emu(5143500)  # 5.625 inches


# ──────────────────────────────────────
# SLIDE 1: TITLE (matching template)
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, TITLE_BG)

# Accent line at top
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04), TITLE_ACCENT)

# "Weekly Reporting"
add_text(slide, Inches(0.75), Inches(1.15), Inches(8.5), Inches(1.0),
         "Weekly Reporting", size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Date
add_text(slide, Inches(1.5), Inches(2.30), Inches(7), Inches(0.5),
         "February 20, 2026", size=24, color=LIGHT_BG, alignment=PP_ALIGN.CENTER)

# Name
add_text(slide, Inches(1.5), Inches(3.10), Inches(7), Inches(0.4),
         "Mohansai Sathram", size=20, color=LIGHT_BG, alignment=PP_ALIGN.CENTER)

# Advisor
add_text(slide, Inches(1.5), Inches(3.70), Inches(7), Inches(0.4),
         "Dr. Yao Hu", size=24, color=LIGHT_BG, alignment=PP_ALIGN.CENTER)

# Bottom accent
add_rect(slide, Inches(0), Inches(5.585), SLIDE_W, Inches(0.04), TITLE_ACCENT)


# ──────────────────────────────────────
# SLIDE 2: LAST WEEK SUMMARY (Part 1)
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Last Week Summary (Feb 13)")

# Card 1: WRF-Hydro Running
add_shape(slide, Inches(0.4), Inches(1.2), Inches(4.4), Inches(1.8), GREEN_LIGHT, GREEN_ACCENT)
add_text(slide, Inches(0.6), Inches(1.3), Inches(4.0), Inches(0.35),
         "WRF-Hydro v5.4.0 Compiled & Running", size=13, color=GREEN_ACCENT, bold=True)
add_bullets(slide, Inches(0.6), Inches(1.65), Inches(4.0), Inches(1.2), [
    "WRF-Hydro = Weather Research & Forecasting\n  Hydrological modeling system by NCAR",
    "Test case: Croton, New York (Hurricane Irene 2011)",
    "6-hour simulation completed successfully (39 output files)",
], size=10, bullet_color=GREEN_ACCENT)

# Card 2: BMI Heat Example
add_shape(slide, Inches(5.1), Inches(1.2), Inches(4.4), Inches(1.8), BLUE_LIGHT, BLUE_ACCENT)
add_text(slide, Inches(5.3), Inches(1.3), Inches(4.0), Inches(0.35),
         "BMI Heat Example: Built & Tested", size=13, color=BLUE_ACCENT, bold=True)
add_bullets(slide, Inches(5.3), Inches(1.65), Inches(4.0), Inches(1.2), [
    "BMI = Basic Model Interface (standard API\n  for model coupling, by CSDMS)",
    "CSDMS = Community Surface Dynamics\n  Modeling System (University of Colorado)",
    "49 out of 49 automated tests passed",
], size=10, bullet_color=BLUE_ACCENT)

# Card 3: WRF-Hydro Internals Studied
add_shape(slide, Inches(0.4), Inches(3.2), Inches(4.4), Inches(2.0), ORANGE_LIGHT, ORANGE_ACCENT)
add_text(slide, Inches(0.6), Inches(3.3), Inches(4.0), Inches(0.35),
         "WRF-Hydro Internals Analyzed", size=13, color=ORANGE_ACCENT, bold=True)
add_bullets(slide, Inches(0.6), Inches(3.65), Inches(4.0), Inches(1.5), [
    "Identified IRF subroutines (Initialize-Run-Finalize)",
    "IRF = the pattern of splitting a model into:\n  initialize(), run_one_step(), finalize()",
    "Mapped 8 key output + 2 input variables\n  to CSDMS Standard Names",
    "CSDMS Standard Names = universal variable\n  naming convention for model coupling",
], size=10, bullet_color=ORANGE_ACCENT)

# Card 4: Documentation
add_shape(slide, Inches(5.1), Inches(3.2), Inches(4.4), Inches(2.0), PURPLE_LIGHT, PURPLE_ACCENT)
add_text(slide, Inches(5.3), Inches(3.3), Inches(4.0), Inches(0.35),
         "Documentation & Planning", size=13, color=PURPLE_ACCENT, bold=True)
add_bullets(slide, Inches(5.3), Inches(3.65), Inches(4.0), Inches(1.5), [
    "Created 8 detailed technical guides covering\n  BMI, CSDMS naming, Babelizer, PyMT",
    "Created Master Implementation Plan\n  (6 phases, 1,115 lines, full roadmap)",
    "Studied SCHISM's BMI approach (how SCHISM\n  enables BMI in its source code)",
    "PyMT = Python Modeling Toolkit (framework\n  that runs coupled models together)",
], size=10, bullet_color=PURPLE_ACCENT)


# ──────────────────────────────────────
# SLIDE 3: THIS WEEK — AGENDA
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "This Week: Key Topics (Feb 20)")

add_text(slide, Inches(0.5), Inches(1.0), Inches(9), Inches(0.35),
         "Three main areas of investigation and progress this week:",
         size=13, color=LIGHT_TEXT, italic=True)

# Topic cards
topics = [
    ("1", "SCHISM's BMI\nApproach", "How SCHISM enables BMI\nusing preprocessor flags\ndirectly in its source code",
     BLUE_ACCENT, BLUE_LIGHT, "Slides 4-6"),
    ("2", "WRF-Hydro:\nWhy Different?", "Why the same approach\nwill not work for WRF-Hydro\nand what we will do instead",
     ORANGE_ACCENT, ORANGE_LIGHT, "Slides 7-9"),
    ("3", "BMI Template &\nHeat Model Example", "The Fortran BMI template\nand how the Heat model\nexample demonstrates it",
     GREEN_ACCENT, GREEN_LIGHT, "Slides 10-14"),
]

for i, (num, title, desc, accent, bg_color, slides) in enumerate(topics):
    x = Inches(0.4 + i * 3.15)
    y = Inches(1.55)

    add_shape(slide, x, y, Inches(2.9), Inches(3.5), bg_color, accent, Pt(1.5))

    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), y + Inches(0.2),
                                   Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.5), Inches(0.7),
             title, size=15, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.2), y + Inches(1.7), Inches(2.5), Inches(1.2),
             desc, size=11, color=MED_TEXT, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.2), y + Inches(3.0), Inches(2.5), Inches(0.3),
             slides, size=10, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────
# SLIDE 4: WHAT IS BMI? (background)
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Background: What is BMI?",
             "BMI = Basic Model Interface — A standardized set of functions that any model must implement to be coupled with other models")

# Left: Explanation
add_shape(slide, Inches(0.4), Inches(1.2), Inches(4.5), Inches(4.1), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.6), Inches(1.3), Inches(4.1), Inches(0.3),
         "BMI in Plain Terms", size=14, color=BLUE_ACCENT, bold=True)

add_bullets(slide, Inches(0.6), Inches(1.7), Inches(4.1), Inches(3.4), [
    "BMI (Basic Model Interface) is a standard API\n  published by CSDMS (University of Colorado Boulder)",
    "It defines 41 functions every model must have,\n  such as: initialize, update, finalize, get_value",
    "Purpose: allow any two models to exchange data\n  without knowing each other's internal code",
    "Think of it like a universal adapter: any model\n  that implements BMI can plug into any other",
    "Written in Fortran, C, C++, or Python — each\n  language has its own BMI specification library",
    "The model's internal code stays untouched;\n  BMI is a wrapper layer on top of it",
], size=10, bullet_color=BLUE_ACCENT, spacing=Pt(5))

# Right: The 6 categories
add_shape(slide, Inches(5.2), Inches(1.2), Inches(4.4), Inches(4.1), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.4), Inches(1.3), Inches(4.0), Inches(0.3),
         "The 41 BMI Functions (6 Categories)", size=14, color=BLUE_ACCENT, bold=True)

add_table(slide, Inches(5.4), Inches(1.7), Inches(4.0), 7, 3, [
    ["Category", "Count", "Key Functions"],
    ["Control", "4", "initialize, update, finalize"],
    ["Model Info", "5", "component_name, var_names"],
    ["Variable Info", "6", "var_type, var_units, var_grid"],
    ["Time", "5", "current_time, time_step"],
    ["Get/Set Data", "6", "get_value, set_value"],
    ["Grid Info", "17", "grid_type, grid_shape, x/y"],
], col_widths=[Inches(1.3), Inches(0.6), Inches(2.1)])

add_text(slide, Inches(5.4), Inches(4.5), Inches(4.0), Inches(0.6),
         "Every function returns a status code:\n0 = BMI_SUCCESS, 1 = BMI_FAILURE\n"
         "Output is always through function parameters,\nnot return values.",
         size=9, color=LIGHT_TEXT, italic=True)


# ──────────────────────────────────────
# SLIDE 5: SCHISM's BMI APPROACH
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "SCHISM's BMI Approach: Inline Preprocessor Flags",
             "SCHISM = Semi-implicit Cross-scale Hydroscience Integrated System Model (ocean/coastal circulation model)")

# Key finding
add_shape(slide, Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.65), ORANGE_LIGHT, ORANGE_ACCENT, Pt(1.5))
add_text(slide, Inches(0.6), Inches(1.25), Inches(8.8), Inches(0.55),
         "Key Finding: SCHISM does NOT have a separate BMI wrapper file. Instead, it enables BMI\n"
         "by adding conditional code blocks (#ifdef USE_NWM_BMI) directly inside its own source files.",
         size=12, color=ORANGE_ACCENT, bold=True, alignment=PP_ALIGN.LEFT)

# Left: What SCHISM does
add_shape(slide, Inches(0.4), Inches(2.05), Inches(4.5), Inches(3.3), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.6), Inches(2.1), Inches(4.1), Inches(0.3),
         "How It Works", size=14, color=BLUE_ACCENT, bold=True)

add_bullets(slide, Inches(0.6), Inches(2.5), Inches(4.1), Inches(2.7), [
    "#ifdef = \"if defined\" — a Fortran/C preprocessor\n"
    "  directive that includes or excludes code at\n"
    "  compile time (before the code is even compiled)",
    "USE_NWM_BMI = a flag name (NWM stands for\n"
    "  National Water Model, NOAA's operational\n"
    "  hydrological model that uses WRF-Hydro)",
    "When flag is OFF: code inside #ifdef is skipped\n"
    "  entirely — SCHISM runs as standalone model",
    "When flag is ON: extra coupling code activates,\n"
    "  allowing SCHISM to receive/send data via BMI",
], size=10, bullet_color=BLUE_ACCENT, spacing=Pt(5))

# Right: Files affected
add_shape(slide, Inches(5.2), Inches(2.05), Inches(4.4), Inches(3.3), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.4), Inches(2.1), Inches(4.0), Inches(0.3),
         "Files with #ifdef USE_NWM_BMI", size=14, color=BLUE_ACCENT, bold=True)

add_table(slide, Inches(5.4), Inches(2.5), Inches(4.0), 5, 3, [
    ["Source File", "Blocks", "Purpose"],
    ["schism_init.F90", "1", "Validates source/sink settings"],
    ["schism_step.F90", "2", "Receives/sends coupling data"],
    ["misc_subs.F90", "2", "Adjusts boundary conditions"],
    ["Total", "5", "Only 5 small code blocks!"],
], col_widths=[Inches(1.5), Inches(0.6), Inches(1.9)])

# Code example
add_text(slide, Inches(5.4), Inches(4.1), Inches(4.0), Inches(0.25),
         "Example from schism_init.F90:", size=10, color=LIGHT_TEXT)
add_code(slide, Inches(5.4), Inches(4.35), Inches(4.0), Inches(0.9),
"""#ifdef USE_NWM_BMI
  if(if_source==0) then
    call parallel_abort('needs if_source>0')
  endif
#endif""", size=9)


# ──────────────────────────────────────
# SLIDE 6: WHY INLINE WORKS FOR SCHISM
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Why This Approach Works for SCHISM",
             "SCHISM's architecture is simple enough for inline preprocessor flags to be manageable")

# 4 reasons in cards
reasons = [
    ("Simple Call Structure",
     "SCHISM's main program calls just 3\nfunctions in sequence:\n\n"
     "  1. schism_init()\n  2. schism_step() [in a loop]\n  3. schism_finalize()\n\n"
     "This is already the Initialize-Run-\n"
     "Finalize (IRF) pattern that BMI needs.\n"
     "No restructuring was required.",
     BLUE_ACCENT, BLUE_LIGHT),
    ("Single Grid Type",
     "SCHISM uses ONE grid throughout:\nan unstructured triangular mesh\n"
     "for the entire ocean domain.\n\n"
     "Every BMI grid function (get_grid_type,\n"
     "get_grid_size, etc.) has exactly\n"
     "one answer — no branching needed.\n\n"
     "This keeps the #ifdef blocks\n"
     "very small and simple.",
     GREEN_ACCENT, GREEN_LIGHT),
    ("Shallow Nesting Depth",
     "The driver (schism_driver.F90) is\nonly 180 lines, 3 levels deep:\n\n"
     "  main -> schism_step -> physics\n\n"
     "#ifdef blocks at level 2 are easy\n"
     "to read and maintain.\n\n"
     "Total files modified: only 3\n"
     "Total #ifdef blocks: only 5",
     TEAL_ACCENT, TEAL_LIGHT),
    ("Coupling Hooks Only",
     "SCHISM's #ifdef blocks are NOT a\nfull 41-function BMI implementation.\n\n"
     "They only add data exchange hooks:\n"
     "  - Receive river discharge input\n"
     "  - Send water elevation output\n\n"
     "This is much simpler than a\n"
     "complete BMI wrapper — just\n"
     "a few lines per #ifdef block.",
     PURPLE_ACCENT, PURPLE_LIGHT),
]

for i, (title, desc, accent, bg_color) in enumerate(reasons):
    x = Inches(0.3 + i * 2.4)
    y = Inches(1.15)
    add_shape(slide, x, y, Inches(2.2), Inches(4.2), bg_color, accent, Pt(1))
    add_text(slide, x + Inches(0.12), y + Inches(0.1), Inches(1.96), Inches(0.45),
             title, size=11, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline(slide, x + Inches(0.12), y + Inches(0.6), Inches(1.96), Inches(3.4),
                  desc, size=9, color=MED_TEXT, line_spacing=Pt(13))


# ──────────────────────────────────────
# SLIDE 7: WRF-HYDRO — WHY DIFFERENT?
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "WRF-Hydro: Can We Use the Same Approach?",
             "WRF-Hydro = Weather Research and Forecasting model - Hydrological extension (by NCAR = National Center for Atmospheric Research)")

# Big answer
add_shape(slide, Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.65), RED_LIGHT, RED_ACCENT, Pt(2))
add_text(slide, Inches(0.6), Inches(1.25), Inches(8.8), Inches(0.55),
         "Answer: No. WRF-Hydro's architecture is fundamentally more complex than SCHISM's.\n"
         "Inline #ifdef flags would require modifying 10+ source files and create unmaintainable code.",
         size=12, color=RED_ACCENT, bold=True)

# Comparison table
add_text(slide, Inches(0.5), Inches(2.1), Inches(9), Inches(0.3),
         "Side-by-Side Comparison:", size=13, color=DARK_TEXT, bold=True)

add_table(slide, Inches(0.5), Inches(2.45), Inches(9), 7, 3, [
    ["Aspect", "SCHISM", "WRF-Hydro"],
    ["Call stack depth", "3 levels (simple)", "5 levels (deeply nested)"],
    ["Main driver size", "180 lines", "2,869 lines"],
    ["Number of grids", "1 (triangular mesh)", "3 (1km + 250m + channel network)"],
    ["Source files", "~50 files", "235 files"],
    ["Time loop structure", "Already init/step/finalize", "Integrated loop, 5 physics phases"],
    ["Physics phases per step", "1 ocean step", "5 sequential phases with dependencies"],
], col_widths=[Inches(2.2), Inches(3.4), Inches(3.4)])

# Bottom insight
add_shape(slide, Inches(0.4), Inches(4.7), Inches(9.2), Inches(0.6), BLUE_LIGHT, BLUE_ACCENT)
add_text(slide, Inches(0.6), Inches(4.75), Inches(8.8), Inches(0.5),
         "Key Difference: SCHISM already had init/step/finalize separation. WRF-Hydro does not — its time\n"
         "loop mixes initialization, forcing, land surface, routing, and output into one large function.",
         size=11, color=BLUE_ACCENT, bold=True)


# ──────────────────────────────────────
# SLIDE 8: WRF-HYDRO 5 PHASES & 3 GRIDS
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "WRF-Hydro: Complex Internal Architecture",
             "Each timestep runs 5 sequential physics phases, and the model uses 3 different computational grids simultaneously")

# Left: 5 Phases
add_shape(slide, Inches(0.4), Inches(1.15), Inches(5.0), Inches(4.2), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.55), Inches(1.2), Inches(4.7), Inches(0.3),
         "5 Physics Phases Per Timestep", size=13, color=BLUE_ACCENT, bold=True)

phases = [
    ("Phase 1: Read Forcing Data", "Reads atmospheric inputs — rainfall, temperature,\nwind, humidity, radiation from input files", BLUE_ACCENT),
    ("Phase 2: NoahMP Land Surface", "Computes soil moisture, evapotranspiration, snow\nprocesses, vegetation dynamics on 1 km grid\n(NoahMP = Noah Multi-Physics land surface model)", GREEN_ACCENT),
    ("Phase 3: Groundwater", "Water table coupling between soil column and\naquifer using WTABLE function", TEAL_ACCENT),
    ("Phase 4: Hydrologic Routing", "Subsurface flow + overland flow on 250 m grid,\nthen channel routing through river network\n(this is HYDRO_exe() — 5 levels deep!)", ORANGE_ACCENT),
    ("Phase 5: Write Output", "Saves results to NetCDF files: streamflow,\nsoil moisture, surface water depth", PURPLE_ACCENT),
]

for i, (name, desc, color) in enumerate(phases):
    y_off = Inches(1.6 + i * 0.72)
    add_text(slide, Inches(0.6), y_off, Inches(1.9), Inches(0.22),
             name, size=9, color=color, bold=True)
    add_text(slide, Inches(2.5), y_off, Inches(2.8), Inches(0.7),
             desc, size=8, color=MED_TEXT)
    if i < 4:
        # Arrow
        add_text(slide, Inches(2.2), y_off + Inches(0.42), Inches(0.3), Inches(0.2),
                 "\u2193", size=10, color=GRAY_BORDER, alignment=PP_ALIGN.CENTER)

# Right: 3 Grids
add_shape(slide, Inches(5.6), Inches(1.15), Inches(4.0), Inches(4.2), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.75), Inches(1.2), Inches(3.7), Inches(0.3),
         "3 Simultaneous Grids", size=13, color=ORANGE_ACCENT, bold=True)

grids = [
    ("Grid 0: LSM Grid (1 km)",
     "Uniform rectilinear grid\nVariables: soil moisture, snow depth,\nevapotranspiration, air temperature\nUsed by: NoahMP land surface model",
     BLUE_ACCENT, BLUE_LIGHT),
    ("Grid 1: Routing Grid (250 m)",
     "Uniform rectilinear grid (4x finer)\nVariables: surface water depth,\nsubsurface flow, overland flow\nUsed by: terrain routing engine",
     GREEN_ACCENT, GREEN_LIGHT),
    ("Grid 2: Channel Network",
     "Unstructured/vector (river reaches)\nVariables: streamflow (discharge),\nchannel coordinates\nUsed by: Muskingum-Cunge routing",
     ORANGE_ACCENT, ORANGE_LIGHT),
]

for i, (name, desc, accent, bg) in enumerate(grids):
    y_off = Inches(1.6 + i * 1.25)
    add_shape(slide, Inches(5.75), y_off, Inches(3.65), Inches(1.1), bg, accent, Pt(0.75))
    add_text(slide, Inches(5.85), y_off + Inches(0.05), Inches(3.45), Inches(0.22),
             name, size=10, color=accent, bold=True)
    add_multiline(slide, Inches(5.85), y_off + Inches(0.3), Inches(3.45), Inches(0.75),
                  desc, size=8, color=MED_TEXT, line_spacing=Pt(11))


# ──────────────────────────────────────
# SLIDE 9: OUR SOLUTION — SEPARATE WRAPPER
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Our Solution: A Separate Non-Invasive BMI Wrapper",
             "We will write bmi_wrf_hydro.f90 — a single Fortran file that wraps WRF-Hydro without modifying any of its source code")

# Architecture diagram as code block
add_text(slide, Inches(0.5), Inches(1.15), Inches(4.5), Inches(0.25),
         "Architecture Overview:", size=12, color=DARK_TEXT, bold=True)

add_code(slide, Inches(0.4), Inches(1.45), Inches(5.5), Inches(3.2),
"""   [External Caller: PyMT / Python / Jupyter]
                     |
          bmi_initialize("config.yaml")
          bmi_update()
          bmi_get_value("streamflow", data)
          bmi_finalize()
                     |
   ========================================
   |    bmi_wrf_hydro.f90  (OUR WRAPPER)  |
   |    - 41 BMI functions                |
   |    - Variable name mapping           |
   |    - Grid info lookup table          |
   ========================================
                     |
         Calls WRF-Hydro subroutines:
         - HYDRO_ini() / HYDRO_exe()
         - land_driver_ini() / _exe()
                     |
   [WRF-Hydro Source Code -- UNTOUCHED]""", size=9)

# Right: Comparison table
add_shape(slide, Inches(6.1), Inches(1.15), Inches(3.6), Inches(3.5), WHITE, GRAY_BORDER)
add_text(slide, Inches(6.25), Inches(1.2), Inches(3.3), Inches(0.3),
         "Why Separate Wrapper Is Better", size=12, color=GREEN_ACCENT, bold=True)

add_table(slide, Inches(6.25), Inches(1.55), Inches(3.3), 7, 2, [
    ["Aspect", "Benefit"],
    ["Non-invasive", "Zero changes to WRF-Hydro"],
    ["Single file", "~800 lines vs 10+ modified files"],
    ["Testable", "Test wrapper independently"],
    ["Maintainable", "Update WRF-Hydro easily"],
    ["Full BMI", "All 41 functions, 3 grids"],
    ["Standard", "Same pattern as BMI examples"],
], col_widths=[Inches(1.2), Inches(2.1)],
   header_bg=GREEN_ACCENT)

# Bottom: The IRF mapping
add_shape(slide, Inches(0.4), Inches(4.8), Inches(9.2), Inches(0.5), BLUE_LIGHT, BLUE_ACCENT)
add_text(slide, Inches(0.6), Inches(4.82), Inches(8.8), Inches(0.45),
         "IRF Mapping:  initialize() \u2192 HYDRO_ini() + land_driver_ini()  |  "
         "update() \u2192 land_driver_exe()  |  finalize() \u2192 HYDRO_finish()",
         size=11, color=BLUE_ACCENT, bold=True)


# ──────────────────────────────────────
# SLIDE 10: BMI FORTRAN TEMPLATE
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "BMI Fortran Template: The Abstract Interface",
             "bmi.f90 — 564 lines that define the contract every BMI-compliant model must fulfill (like an interface specification)")

# Left: What it is
add_shape(slide, Inches(0.4), Inches(1.2), Inches(4.5), Inches(2.2), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.6), Inches(1.3), Inches(4.1), Inches(0.3),
         "What is bmi.f90?", size=14, color=BLUE_ACCENT, bold=True)
add_bullets(slide, Inches(0.6), Inches(1.65), Inches(4.1), Inches(1.6), [
    "An abstract type definition in Fortran 2003 — it declares\n  53 procedures but provides NO implementation",
    "Any model that wants BMI must \"extend\" this type\n  and provide its own implementation of each function",
    "Published by CSDMS as the bmif_2_0 module\n  (BMI version 2.0, Hutton et al. 2020)",
    "Installed via conda: bmi-fortran package provides\n  the compiled module and library (libbmif.so)",
], size=10, bullet_color=BLUE_ACCENT, spacing=Pt(4))

# Right: Structure
add_shape(slide, Inches(5.2), Inches(1.2), Inches(4.4), Inches(2.2), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.4), Inches(1.3), Inches(4.0), Inches(0.3),
         "Key Constants Defined in bmi.f90", size=14, color=BLUE_ACCENT, bold=True)

add_code(slide, Inches(5.4), Inches(1.7), Inches(4.0), Inches(1.5),
"""! Status codes returned by every BMI function:
integer, parameter :: BMI_SUCCESS = 0
integer, parameter :: BMI_FAILURE = 1

! Maximum string lengths for names:
integer, parameter :: BMI_MAX_COMPONENT_NAME = 2048
integer, parameter :: BMI_MAX_VAR_NAME = 2048
integer, parameter :: BMI_MAX_TYPE_NAME = 2048
integer, parameter :: BMI_MAX_UNITS_NAME = 2048""", size=9)

# Bottom: How it works
add_shape(slide, Inches(0.4), Inches(3.6), Inches(9.2), Inches(1.7), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.6), Inches(3.65), Inches(8.8), Inches(0.3),
         "How Extending Works (Type Inheritance in Fortran):", size=13, color=BLUE_ACCENT, bold=True)

add_code(slide, Inches(0.5), Inches(4.0), Inches(4.4), Inches(1.2),
"""! bmi.f90 declares (abstract - no code):
type, abstract :: bmi
contains
  procedure(bmif_initialize), deferred :: initialize
  procedure(bmif_update), deferred :: update
  procedure(bmif_finalize), deferred :: finalize
  ! ... 50 more deferred procedures
end type""", size=9)

add_code(slide, Inches(5.1), Inches(4.0), Inches(4.4), Inches(1.2),
"""! Our wrapper EXTENDS it (provides code):
type, extends(bmi) :: bmi_wrf_hydro
  private
  ! embed WRF-Hydro model state here
contains
  procedure :: initialize => wrfhydro_init
  procedure :: update => wrfhydro_update
  procedure :: finalize => wrfhydro_finalize
  ! ... implement all 53 procedures
end type""", size=9)


# ──────────────────────────────────────
# SLIDE 11: HEAT MODEL — THE PHYSICS
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "BMI Heat Model Example: The Physics Model",
             "heat.f90 — A simple 2D heat diffusion model (158 lines). This is the model that gets wrapped by BMI, just like WRF-Hydro will be.")

# Left: The model
add_shape(slide, Inches(0.4), Inches(1.2), Inches(4.5), Inches(4.1), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.6), Inches(1.3), Inches(4.1), Inches(0.3),
         "What Does the Heat Model Do?", size=14, color=BLUE_ACCENT, bold=True)

add_bullets(slide, Inches(0.6), Inches(1.7), Inches(4.1), Inches(1.5), [
    "Simulates heat spreading across a 2D metal plate",
    "Top edge is held at a fixed hot temperature (like\n  a boundary condition in hydrology)",
    "Heat gradually diffuses to neighboring cells using\n  the 2D heat equation (finite difference method)",
    "Each timestep: new_temp = weighted average of 4\n  neighboring cells (a 5-point stencil computation)",
], size=10, bullet_color=BLUE_ACCENT, spacing=Pt(5))

add_text(slide, Inches(0.6), Inches(3.2), Inches(4.1), Inches(0.25),
         "Model State Variables (stored in heat_model type):", size=11, color=DARK_TEXT, bold=True)
add_code(slide, Inches(0.5), Inches(3.5), Inches(4.3), Inches(1.7),
"""type :: heat_model
  integer :: n_x = 0    ! Grid columns (width)
  integer :: n_y = 0    ! Grid rows (height)
  real :: dt = 0.       ! Time step (seconds)
  real :: t  = 0.       ! Current simulation time
  real :: t_end = 0.    ! End time of simulation
  real :: alpha = 0.    ! Thermal diffusivity

  ! Main state variable (2D temperature field):
  real, pointer :: temperature(:,:) => null()
end type""", size=9)

# Right: 3 subroutines
add_shape(slide, Inches(5.2), Inches(1.2), Inches(4.4), Inches(4.1), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.4), Inches(1.3), Inches(4.0), Inches(0.3),
         "3 Model Subroutines", size=14, color=GREEN_ACCENT, bold=True)

subs = [
    ("initialize_from_file(model, filename)",
     "Reads 4 values from config file:\nalpha (diffusivity), t_end (end time),\nn_x (columns), n_y (rows).\nAllocates the temperature grid.\nSets initial condition (top row = hot).",
     BLUE_ACCENT, BLUE_LIGHT),
    ("advance_in_time(model)",
     "Runs ONE timestep of heat diffusion:\nFor each interior cell, computes new\ntemperature as weighted average of its\n4 neighbors. Copies result back.\nIncrements time by dt.",
     GREEN_ACCENT, GREEN_LIGHT),
    ("cleanup(model)",
     "Deallocates the temperature arrays.\nFrees all memory used by the model.\nCalled at the end of the simulation.",
     RED_ACCENT, RED_LIGHT),
]

for i, (name, desc, accent, bg) in enumerate(subs):
    y_off = Inches(1.7 + i * 1.2)
    add_shape(slide, Inches(5.35), y_off, Inches(4.1), Inches(1.05), bg, accent, Pt(0.75))
    add_text(slide, Inches(5.45), y_off + Inches(0.03), Inches(3.9), Inches(0.2),
             name, size=9, color=accent, bold=True, font=CODE_FONT)
    add_multiline(slide, Inches(5.45), y_off + Inches(0.25), Inches(3.9), Inches(0.75),
                  desc, size=8, color=MED_TEXT, line_spacing=Pt(11))

# Analogy box
add_shape(slide, Inches(5.35), Inches(4.35), Inches(4.1), Inches(0.85), TEAL_LIGHT, TEAL_ACCENT)
add_text(slide, Inches(5.45), Inches(4.38), Inches(3.9), Inches(0.2),
         "Hydrology Analogy:", size=10, color=TEAL_ACCENT, bold=True)
add_text(slide, Inches(5.45), Inches(4.6), Inches(3.9), Inches(0.55),
         "This simple model is to bmi_heat.f90 what WRF-Hydro\n"
         "is to our bmi_wrf_hydro.f90 — the underlying physics\n"
         "engine that the BMI wrapper calls into.",
         size=9, color=MED_TEXT)


# ──────────────────────────────────────
# SLIDE 12: BMI HEAT WRAPPER — PATTERNS
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "bmi_heat.f90: The BMI Wrapper (935 Lines)",
             "This file is our blueprint. Every pattern here will be reused when writing bmi_wrf_hydro.f90.")

# Pattern 1: Type Extension
add_shape(slide, Inches(0.4), Inches(1.2), Inches(4.5), Inches(2.0), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.55), Inches(1.25), Inches(4.2), Inches(0.25),
         "Pattern 1: Type Extension (Embed the Model)", size=12, color=BLUE_ACCENT, bold=True)
add_code(slide, Inches(0.5), Inches(1.55), Inches(4.3), Inches(1.55),
"""module bmiheatf
  use heatf          ! Import heat model module
  use bmif_2_0       ! Import BMI specification
  use, intrinsic :: iso_c_binding  ! For pointers

  type, extends(bmi) :: bmi_heat
    private
    type(heat_model) :: model  ! Embed model instance
  contains
    procedure :: initialize => heat_initialize
    procedure :: update => heat_update
    ! ... maps all 53 procedures
  end type bmi_heat""", size=9)

# Pattern 2: Select Case
add_shape(slide, Inches(5.2), Inches(1.2), Inches(4.4), Inches(2.0), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.35), Inches(1.25), Inches(4.1), Inches(0.25),
         "Pattern 2: Variable Name Dispatch", size=12, color=GREEN_ACCENT, bold=True)
add_code(slide, Inches(5.3), Inches(1.55), Inches(4.2), Inches(1.55),
"""! Every variable-info function uses this:
function heat_var_type(self, name, type) ...
  select case(name)
  case("plate_surface__temperature")
    type = "real"
    bmi_status = BMI_SUCCESS
  case default
    type = "-"
    bmi_status = BMI_FAILURE  ! Unknown var
  end select
end function""", size=9)

# Pattern 3: get_value with reshape
add_shape(slide, Inches(0.4), Inches(3.4), Inches(4.5), Inches(1.9), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.55), Inches(3.45), Inches(4.2), Inches(0.25),
         "Pattern 3: get_value (Copy + Flatten to 1D)", size=12, color=ORANGE_ACCENT, bold=True)
add_code(slide, Inches(0.5), Inches(3.75), Inches(4.3), Inches(1.45),
"""! BMI requires all arrays as 1D (flattened):
function heat_get_double(self, name, dest) ...
  select case(name)
  case("plate_surface__temperature")
    ! Flatten 2D grid -> 1D array:
    dest = reshape(self%model%temperature, &
                   [self%model%n_y * self%model%n_x])
    bmi_status = BMI_SUCCESS
  end select
end function""", size=9)

# Pattern 4: Control functions
add_shape(slide, Inches(5.2), Inches(3.4), Inches(4.4), Inches(1.9), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.35), Inches(3.45), Inches(4.1), Inches(0.25),
         "Pattern 4: Control Functions (IRF)", size=12, color=PURPLE_ACCENT, bold=True)
add_code(slide, Inches(5.3), Inches(3.75), Inches(4.2), Inches(1.45),
"""! Initialize: delegate to model's own init
function heat_initialize(self, config_file) ...
  call initialize_from_file(self%model, &
                            config_file)
end function

! Update: advance by exactly ONE timestep
function heat_update(self) ...
  call advance_in_time(self%model)
end function""", size=9)


# ──────────────────────────────────────
# SLIDE 13: BMI DRIVER & DATA EXCHANGE
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "BMI Driver Program & Data Exchange",
             "bmi_main.f90 (65 lines) — Shows how an external program uses BMI to control the model. The caller owns the time loop, not the model.")

# Left: Driver program
add_shape(slide, Inches(0.4), Inches(1.2), Inches(4.5), Inches(4.1), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.6), Inches(1.25), Inches(4.1), Inches(0.25),
         "BMI Driver: The Caller Controls Everything", size=12, color=BLUE_ACCENT, bold=True)

add_code(slide, Inches(0.5), Inches(1.55), Inches(4.3), Inches(3.6),
"""program bmi_main
  use bmiheatf          ! Import the BMI wrapper
  type(bmi_heat) :: m   ! Create model instance

  ! Step 1: INITIALIZE (read config, allocate)
  status = m%initialize("config.cfg")

  ! Step 2: QUERY MODEL INFO
  status = m%get_component_name(name)
  status = m%get_start_time(time)
  status = m%get_end_time(end_time)
  status = m%get_time_step(dt)

  ! Step 3: TIME LOOP (caller controls it!)
  do while (current_time < end_time)
    ! Read data OUT of the model:
    status = m%get_value( &
        "plate_surface__temperature", temp)
    print *, "Temperature:", temp(1:5)

    ! Advance model by one timestep:
    status = m%update()

    status = m%get_current_time(current_time)
  end do

  ! Step 4: CLEAN UP
  status = m%finalize()
end program""", size=8)

# Right: How coupling works
add_shape(slide, Inches(5.2), Inches(1.2), Inches(4.4), Inches(2.2), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.4), Inches(1.25), Inches(4.0), Inches(0.25),
         "How Two Models Exchange Data via BMI", size=12, color=GREEN_ACCENT, bold=True)

add_code(slide, Inches(5.3), Inches(1.55), Inches(4.2), Inches(1.75),
"""! Future: WRF-Hydro + SCHISM coupling loop
do while (time < end_time)
  ! Run WRF-Hydro for one timestep
  call wrfhydro%update()

  ! Get river discharge from WRF-Hydro
  call wrfhydro%get_value( &
      "channel_water__volume_flow_rate", Q)

  ! Feed discharge into SCHISM
  call schism%set_value( &
      "channel_water__volume_flow_rate", Q)

  ! Run SCHISM for one timestep
  call schism%update()
end do""", size=9)

# Right bottom: Key points
add_shape(slide, Inches(5.2), Inches(3.6), Inches(4.4), Inches(1.7), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.4), Inches(3.65), Inches(4.0), Inches(0.25),
         "Key Observations", size=12, color=ORANGE_ACCENT, bold=True)

add_bullets(slide, Inches(5.4), Inches(3.95), Inches(4.0), Inches(1.3), [
    "The caller owns the time loop — models do not\n  run themselves; they are stepped externally",
    "All data access is through get_value/set_value\n  using CSDMS Standard Names (not internal names)",
    "Config file is the only argument to initialize() —\n  all model settings are read from that file",
    "Every BMI call returns a status code (0 or 1) —\n  production code should always check these",
], size=9, bullet_color=ORANGE_ACCENT, spacing=Pt(3))


# ──────────────────────────────────────
# SLIDE 14: HOW WE'LL APPLY THIS TO WRF-HYDRO
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Applying the BMI Heat Pattern to WRF-Hydro",
             "Mapping every pattern from bmi_heat.f90 to our bmi_wrf_hydro.f90 — same structure, different model")

# Mapping table
add_text(slide, Inches(0.5), Inches(1.15), Inches(9), Inches(0.3),
         "Pattern-by-Pattern Mapping:", size=13, color=DARK_TEXT, bold=True)

add_table(slide, Inches(0.4), Inches(1.5), Inches(9.2), 9, 3, [
    ["BMI Pattern", "Heat Model (Template)", "WRF-Hydro (Our Wrapper)"],
    ["Module name", "bmiheatf", "bmiwrfhydrof"],
    ["Embedded model", "type(heat_model) :: model", "WRF-Hydro state variables (RT_FIELD)"],
    ["initialize()", "initialize_from_file()", "HYDRO_ini() + land_driver_ini()"],
    ["update()", "advance_in_time()", "land_driver_exe() (all 5 phases)"],
    ["finalize()", "cleanup()", "HYDRO_finish() + close files"],
    ["Grid 0", "1 uniform_rectilinear (nx * ny)", "1 km LSM grid (ix * jx)"],
    ["Variables", "temperature (1 variable)", "streamflow, soil moisture, snow... (8+)"],
    ["Name mapping", "plate_surface__temperature", "channel_water__volume_flow_rate, etc."],
], col_widths=[Inches(1.8), Inches(3.7), Inches(3.7)])

# Bottom: variable mapping example
add_text(slide, Inches(0.5), Inches(4.45), Inches(9), Inches(0.25),
         "CSDMS Standard Names Mapping for WRF-Hydro (our initial 8 variables):", size=11, color=DARK_TEXT, bold=True)

add_table(slide, Inches(0.4), Inches(4.75), Inches(9.2), 5, 4, [
    ["Internal Name", "CSDMS Standard Name", "Units", "Grid"],
    ["QLINK(:,1)", "channel_water__volume_flow_rate", "m\u00B3/s", "Channel (#2)"],
    ["sfcheadrt", "land_surface_water__depth", "m", "Routing (#1)"],
    ["SOIL_M", "soil_water__volume_fraction", "dimensionless", "LSM (#0)"],
    ["SNEQV", "snowpack__liquid-equivalent_depth", "m", "LSM (#0)"],
], col_widths=[Inches(1.5), Inches(3.7), Inches(1.5), Inches(2.5)])


# ──────────────────────────────────────
# SLIDE 15: NEXT STEPS
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Next Steps: Implementation Roadmap")

# Completed
add_shape(slide, Inches(0.4), Inches(1.1), Inches(4.5), Inches(4.2), GREEN_LIGHT, GREEN_ACCENT, Pt(1))
add_text(slide, Inches(0.6), Inches(1.15), Inches(4.1), Inches(0.35),
         "Completed (Phase 1a)", size=14, color=GREEN_ACCENT, bold=True)

done = [
    "Studied BMI specification (bmi.f90, 41 functions)",
    "Built and tested BMI Heat example (49/49 tests pass)",
    "Compiled WRF-Hydro v5.4.0 successfully",
    "Ran Croton, NY test case (6-hour Hurricane Irene)",
    "Analyzed WRF-Hydro call stack and IRF subroutines",
    "Studied SCHISM's BMI approach (#ifdef pattern)",
    "Mapped 8 WRF-Hydro variables to CSDMS names",
    "Created Master Implementation Plan (1,115 lines)",
    "Created 8 detailed technical reference documents",
]
add_bullets(slide, Inches(0.55), Inches(1.5), Inches(4.2), Inches(3.6), done,
            size=10, bullet_char="\u2713", bullet_color=GREEN_ACCENT, spacing=Pt(3))

# Upcoming
add_shape(slide, Inches(5.2), Inches(1.1), Inches(4.4), Inches(2.4), ORANGE_LIGHT, ORANGE_ACCENT, Pt(1))
add_text(slide, Inches(5.4), Inches(1.15), Inches(4.0), Inches(0.35),
         "Immediate Next (Phase 1b)", size=14, color=ORANGE_ACCENT, bold=True)

upcoming = [
    "Refactor WRF-Hydro time loop into IRF pattern\n  (separate init, one-step, finalize)",
    "Write bmi_wrf_hydro.f90 with all 41 functions",
    "Compile as shared library (libwrfhydro_bmi.so)",
    "Write Fortran test driver to validate BMI output\n  against standalone WRF-Hydro run",
]
add_bullets(slide, Inches(5.35), Inches(1.5), Inches(4.1), Inches(1.8), upcoming,
            size=10, bullet_char="\u25B6", bullet_color=ORANGE_ACCENT, spacing=Pt(4))

# Future
add_shape(slide, Inches(5.2), Inches(3.65), Inches(4.4), Inches(1.65), PURPLE_LIGHT, PURPLE_ACCENT, Pt(1))
add_text(slide, Inches(5.4), Inches(3.7), Inches(4.0), Inches(0.35),
         "Future Phases (2-4)", size=14, color=PURPLE_ACCENT, bold=True)

future = [
    "Babelize both models (Fortran \u2192 Python plugins)",
    "Register as PyMT plugins for coupling framework",
    "Write coupling script (WRF-Hydro \u2194 SCHISM)",
    "Test compound flooding case study",
]
add_bullets(slide, Inches(5.35), Inches(4.05), Inches(4.1), Inches(1.2), future,
            size=10, bullet_char="\u25CB", bullet_color=PURPLE_ACCENT, spacing=Pt(3))


# ──────────────────────────────────────
# SLIDE 16: QLINK — WHERE STREAMFLOW LIVES
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "QLINK: The Key Coupling Variable (Streamflow)",
             "QLINK = the Fortran variable inside WRF-Hydro that stores river discharge (streamflow) at every channel link")

# Top banner: CSDMS name
add_shape(slide, Inches(0.4), Inches(1.15), Inches(9.2), Inches(0.55), GREEN_LIGHT, GREEN_ACCENT, Pt(1.5))
add_text(slide, Inches(0.6), Inches(1.18), Inches(8.8), Inches(0.5),
         "CSDMS Standard Name:  channel_water__volume_flow_rate    |    Units: m\u00B3/s    |    Grid: Channel Network",
         size=13, color=GREEN_ACCENT, bold=True)

# Left: Where QLINK lives
add_shape(slide, Inches(0.4), Inches(1.9), Inches(4.5), Inches(3.4), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.55), Inches(1.95), Inches(4.2), Inches(0.3),
         "Where QLINK Lives in Source Code", size=13, color=BLUE_ACCENT, bold=True)

add_table(slide, Inches(0.5), Inches(2.35), Inches(4.3), 6, 2, [
    ["What", "Details"],
    ["Declaration", "module_rt_inc.F90 (line 160)\nInside RT_FIELD derived type"],
    ["Data type", "REAL, allocatable, DIMENSION(:,:)\nQLINK(nlinks, 2)"],
    ["Column 1", "QLINK(:,1) = previous timestep flow"],
    ["Column 2", "QLINK(:,2) = current timestep flow"],
    ["Access at runtime", "RT_DOMAIN(1)%QLINK(1:nlinks, 2)"],
], col_widths=[Inches(1.3), Inches(3.0)],
   header_bg=BLUE_ACCENT)

# Right: Data flow
add_shape(slide, Inches(5.2), Inches(1.9), Inches(4.4), Inches(1.8), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.35), Inches(1.95), Inches(4.1), Inches(0.3),
         "How QLINK Flows Through WRF-Hydro", size=12, color=ORANGE_ACCENT, bold=True)

add_code(slide, Inches(5.3), Inches(2.3), Inches(4.2), Inches(1.3),
"""1. Forcing data (rain) enters WRF-Hydro
2. NoahMP computes soil water & runoff
3. Overland routing moves water to channels
4. Muskingum-Cunge routing computes:
     QLINK(k,2) = flow at link k (m3/s)
5. QLINK written to CHRTOUT output files
     as "streamflow" variable in NetCDF""", size=9)

# Right bottom: Output file info
add_shape(slide, Inches(5.2), Inches(3.85), Inches(4.4), Inches(1.45), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.35), Inches(3.9), Inches(4.1), Inches(0.3),
         "Files That Compute / Store QLINK", size=12, color=PURPLE_ACCENT, bold=True)

add_table(slide, Inches(5.3), Inches(4.2), Inches(4.2), 5, 2, [
    ["Source File", "Role"],
    ["module_rt_inc.F90", "Declares QLINK in RT_FIELD type"],
    ["module_channel_routing.F90", "Computes QLINK via routing"],
    ["module_HYDRO_io.F90", "Writes QLINK to CHRTOUT files"],
    ["module_HYDRO_drv.F90", "Passes QLINK to output routine"],
], col_widths=[Inches(2.2), Inches(2.0)],
   header_bg=PURPLE_ACCENT)


# ──────────────────────────────────────
# SLIDE 17: WRF-HYDRO → SCHISM COUPLING VIA QLINK
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Coupling: WRF-Hydro QLINK \u2192 SCHISM ath3 Array",
             "How river discharge computed by WRF-Hydro enters SCHISM's coastal ocean model as source/sink boundary conditions")

# Architecture flow
add_text(slide, Inches(0.5), Inches(1.15), Inches(9), Inches(0.25),
         "Data Flow: WRF-Hydro Streamflow \u2192 SCHISM Ocean Boundary", size=12, color=DARK_TEXT, bold=True)

add_code(slide, Inches(0.4), Inches(1.45), Inches(9.2), Inches(1.6),
"""   WRF-HYDRO                          BMI LAYER                          SCHISM
   ============                    ==================                   ============
   module_channel_routing.F90       bmi_wrf_hydro.f90                   schism_step.F90
   QLINK(k,2) computed by     -->  get_value("channel_water            ath3(:,1,2,1) receives
   Muskingum-Cunge routing          __volume_flow_rate", Q)             river discharge at
   for each river link              returns QLINK as 1D array   -->    source/sink nodes
                                                                        (via #ifdef USE_NWM_BMI)
   Variable: RT_DOMAIN(1)%QLINK     CSDMS Standard Name:               Variable: ath3(nsources,
   Shape: (nlinks, 2)               channel_water__volume_flow_rate     1, 2, 1)
   Units: m3/s                      Flatten: QLINK(:,2) -> 1D           Units: m3/s""", size=8)

# Bottom left: SCHISM receives
add_shape(slide, Inches(0.4), Inches(3.2), Inches(4.5), Inches(2.1), BLUE_LIGHT, BLUE_ACCENT, Pt(1))
add_text(slide, Inches(0.55), Inches(3.25), Inches(4.2), Inches(0.3),
         "How SCHISM Receives Discharge (ath3 array)", size=12, color=BLUE_ACCENT, bold=True)

add_bullets(slide, Inches(0.55), Inches(3.6), Inches(4.2), Inches(1.6), [
    "ath3 = SCHISM's source/sink input array\n  (declared in schism_glbl.F90, line 386)",
    "ath3(nsources, 1, 1:2, 1) stores flow rates at\n  two time levels (for interpolation)",
    "When USE_NWM_BMI is ON, SCHISM skips reading\n  ath3 from files — BMI provides it directly",
    "Location: schism_step.F90 (line 1540) and\n  misc_subs.F90 (line 599) #ifdef blocks",
], size=10, bullet_color=BLUE_ACCENT, spacing=Pt(4))

# Bottom right: Current status
add_shape(slide, Inches(5.2), Inches(3.2), Inches(4.4), Inches(2.1), ORANGE_LIGHT, ORANGE_ACCENT, Pt(1))
add_text(slide, Inches(5.35), Inches(3.25), Inches(4.1), Inches(0.3),
         "Current Coupling Support Status", size=12, color=ORANGE_ACCENT, bold=True)

add_table(slide, Inches(5.35), Inches(3.6), Inches(4.1), 4, 3, [
    ["Direction", "Variable", "Status"],
    ["WRF-Hydro \u2192 SCHISM", "Discharge (QLINK)", "Supported (ath3)"],
    ["SCHISM \u2192 WRF-Hydro", "Water Elevation (eta2)", "Needs Addition"],
    ["Both directions", "Full 2-way coupling", "Phase 4 Goal"],
], col_widths=[Inches(1.5), Inches(1.5), Inches(1.1)],
   header_bg=ORANGE_ACCENT)


# ──────────────────────────────────────
# SLIDE 18: QUESTIONS FOR PROFESSOR
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Questions for Discussion",
             "Key decisions needed before proceeding with BMI wrapper implementation")

# Question 1: Which variables
add_shape(slide, Inches(0.4), Inches(1.15), Inches(4.5), Inches(2.0), BLUE_LIGHT, BLUE_ACCENT, Pt(1.5))
add_text(slide, Inches(0.55), Inches(1.2), Inches(4.2), Inches(0.3),
         "Q1: Which WRF-Hydro Variables for BMI?", size=13, color=BLUE_ACCENT, bold=True)
add_bullets(slide, Inches(0.55), Inches(1.55), Inches(4.2), Inches(1.5), [
    "WRF-Hydro has 150+ output variables across\n  8 output file types (CHRTOUT, LDASOUT, etc.)",
    "BMI does NOT require exposing all variables —\n  we can start with a priority subset",
    "Currently proposed: 8 output + 2 input variables\n  (streamflow, soil moisture, snow, precip, etc.)",
    "Which variables are most important for\n  our coupling use case with SCHISM?",
], size=10, bullet_color=BLUE_ACCENT, spacing=Pt(4))

# Question 2: SCHISM BMI support
add_shape(slide, Inches(5.2), Inches(1.15), Inches(4.4), Inches(2.0), GREEN_LIGHT, GREEN_ACCENT, Pt(1.5))
add_text(slide, Inches(5.35), Inches(1.2), Inches(4.1), Inches(0.3),
         "Q2: Does SCHISM BMI Support These?", size=13, color=GREEN_ACCENT, bold=True)
add_bullets(slide, Inches(5.35), Inches(1.55), Inches(4.1), Inches(1.5), [
    "SCHISM's #ifdef USE_NWM_BMI currently only\n  supports receiving discharge (via ath3 array)",
    "It does NOT export water elevation (eta2)\n  back to WRF-Hydro — needed for 2-way coupling",
    "For 1-way coupling (WRF-Hydro \u2192 SCHISM):\n  QLINK (streamflow) is sufficient",
    "For 2-way: SCHISM needs modifications to\n  export coastal water level as BMI variable",
], size=10, bullet_color=GREEN_ACCENT, spacing=Pt(4))

# Question 3: Scope of initial implementation
add_shape(slide, Inches(0.4), Inches(3.35), Inches(4.5), Inches(1.9), ORANGE_LIGHT, ORANGE_ACCENT, Pt(1.5))
add_text(slide, Inches(0.55), Inches(3.4), Inches(4.2), Inches(0.3),
         "Q3: Start with 1-Way or 2-Way Coupling?", size=13, color=ORANGE_ACCENT, bold=True)
add_bullets(slide, Inches(0.55), Inches(3.75), Inches(4.2), Inches(1.4), [
    "1-Way: WRF-Hydro sends streamflow to SCHISM\n  (simpler, SCHISM already supports receiving it)",
    "2-Way: Also get coastal water elevation back\n  (requires modifying SCHISM's BMI blocks)",
    "Recommendation: Start 1-way, add 2-way later",
    "Which approach should we prioritize?",
], size=10, bullet_color=ORANGE_ACCENT, spacing=Pt(4))

# Question 4: CSDMS Names
add_shape(slide, Inches(5.2), Inches(3.35), Inches(4.4), Inches(1.9), PURPLE_LIGHT, PURPLE_ACCENT, Pt(1.5))
add_text(slide, Inches(5.35), Inches(3.4), Inches(4.1), Inches(0.3),
         "Q4: CSDMS Standard Names Registration?", size=13, color=PURPLE_ACCENT, bold=True)
add_bullets(slide, Inches(5.35), Inches(3.75), Inches(4.1), Inches(1.4), [
    "Only ~22% of WRF-Hydro variables have\n  existing CSDMS Standard Names",
    "~78% need new names to be proposed\n  and registered with CSDMS",
    "Should we propose new names now, or use\n  temporary names and register later?",
    "Registration process: submit to CSDMS wiki\n  for community review",
], size=10, bullet_color=PURPLE_ACCENT, spacing=Pt(4))


# ──────────────────────────────────────
# SLIDE 19: THANK YOU
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)

# Center content
add_text(slide, Inches(1), Inches(1.5), Inches(8), Inches(0.6),
         "Thank You", size=36, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(2.3), Inches(8), Inches(0.4),
         "Questions & Discussion", size=20, color=BLUE_ACCENT, alignment=PP_ALIGN.CENTER)

# Summary cards
cards = [
    ("SCHISM", "Inline #ifdef works because\nof simple architecture:\nsingle grid, shallow nesting,\nalready has IRF pattern", BLUE_ACCENT, BLUE_LIGHT),
    ("WRF-Hydro", "Needs separate wrapper due\nto 5-level nesting, 3 grids,\n5 physics phases per step,\n235 source files", ORANGE_ACCENT, ORANGE_LIGHT),
    ("Our Approach", "Write bmi_wrf_hydro.f90:\nnon-invasive, single file,\nfollowing bmi_heat.f90\ntemplate patterns", GREEN_ACCENT, GREEN_LIGHT),
]

for i, (title, desc, accent, bg) in enumerate(cards):
    x = Inches(0.7 + i * 3.1)
    y = Inches(3.2)
    add_shape(slide, x, y, Inches(2.8), Inches(1.9), bg, accent, Pt(1.5))
    add_text(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.5), Inches(0.35),
             title, size=14, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.5), Inches(1.2),
                  desc, size=10, color=MED_TEXT, line_spacing=Pt(14))

# Bottom
add_rect(slide, Inches(0), Inches(5.3), SLIDE_W, Inches(0.325), BLUE_ACCENT)
add_text(slide, Inches(1), Inches(5.32), Inches(8), Inches(0.3),
         "WRF-Hydro BMI Wrapper Project  |  Mohansai Sathram  |  February 20, 2026",
         size=10, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════
output_path = "/mnt/c/Users/mohansai/Desktop/Projects/VS_Code/WRF-Hydro-BMI/bmi_wrf_hydro/Docs/Weekly Reporting/20_Feb_2026_WRF_Hydro_BMI_v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
