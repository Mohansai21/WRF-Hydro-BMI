#!/usr/bin/env python3
"""
Generate SCHISM BMI Deep Dive Presentation
Light theme, Times New Roman, detailed technical content
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette (Light Theme) ───
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT_BG = RGBColor(0xF0, 0xF2, 0xF5)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MED_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_TEXT = RGBColor(0x55, 0x55, 0x55)
BLUE_ACCENT = RGBColor(0x1A, 0x56, 0xDB)
BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)
BLUE_MED = RGBColor(0xD0, 0xE1, 0xFD)
GREEN_ACCENT = RGBColor(0x0D, 0x7C, 0x3E)
GREEN_LIGHT = RGBColor(0xE6, 0xF4, 0xEA)
ORANGE_ACCENT = RGBColor(0xC6, 0x5D, 0x07)
ORANGE_LIGHT = RGBColor(0xFE, 0xF3, 0xE2)
RED_ACCENT = RGBColor(0xC5, 0x22, 0x1F)
RED_LIGHT = RGBColor(0xFC, 0xE8, 0xE6)
PURPLE_ACCENT = RGBColor(0x6A, 0x1B, 0x9A)
PURPLE_LIGHT = RGBColor(0xF3, 0xE5, 0xF5)
TEAL_ACCENT = RGBColor(0x00, 0x69, 0x6B)
TEAL_LIGHT = RGBColor(0xE0, 0xF2, 0xF1)
GRAY_BORDER = RGBColor(0xDA, 0xDC, 0xE0)
DARK_CARD = RGBColor(0x26, 0x32, 0x38)
CODE_GREEN = RGBColor(0x4C, 0xAF, 0x50)
TITLE_BG = RGBColor(0x1A, 0x1A, 0x2E)
TITLE_ACCENT = RGBColor(0x3D, 0x8B, 0xFD)

FONT = "Times New Roman"
CODE_FONT = "Consolas"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0.75)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.03
    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0.75)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=14, color=DARK_TEXT,
             bold=False, alignment=PP_ALIGN.LEFT, font=FONT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, text, size=12, color=MED_TEXT,
                  bold=False, font=FONT, line_spacing=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.space_after = Pt(2)
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def add_bullets(slide, left, top, width, height, items, size=12, color=MED_TEXT,
                bullet_char="\u2022", bullet_color=BLUE_ACCENT, spacing=Pt(4)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.space_before = Pt(1)
        run_b = p.add_run()
        run_b.text = f"{bullet_char} "
        run_b.font.size = Pt(size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = FONT
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(size)
        run_t.font.color.rgb = color
        run_t.font.name = FONT
    return txBox


def add_code(slide, left, top, width, height, code_text, size=9):
    shape = add_shape(slide, left, top, width, height, DARK_CARD, GRAY_BORDER, Pt(0.5))
    txBox = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.08),
                                     width - Inches(0.24), height - Inches(0.16))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = CODE_GREEN
        p.font.name = CODE_FONT
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    return shape


def slide_header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.72), BLUE_ACCENT)
    add_text(slide, Inches(0.5), Inches(0.12), Inches(9), Inches(0.5),
             title, size=22, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.72), Inches(9), Inches(0.35),
                 subtitle, size=11, color=LIGHT_TEXT, italic=True)


def add_table(slide, left, top, width, rows, cols, data, col_widths=None,
              header_bg=BLUE_ACCENT, header_fg=WHITE, row_bg1=WHITE, row_bg2=LIGHT_BG):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.32 * rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = header_fg
                else:
                    p.font.color.rgb = MED_TEXT
            cf = cell.fill
            cf.solid()
            if r == 0:
                cf.fore_color.rgb = header_bg
            elif r % 2 == 0:
                cf.fore_color.rgb = row_bg2
            else:
                cf.fore_color.rgb = row_bg1
    return table_shape


# ═══════════════════════════════════════
#  BUILD PRESENTATION
# ═══════════════════════════════════════

prs = Presentation()
prs.slide_width = Emu(9144000)   # 10 inches
prs.slide_height = Emu(5143500)  # 5.625 inches


# ──────────────────────────────────────
# SLIDE 1: TITLE
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, TITLE_BG)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04), TITLE_ACCENT)

add_text(slide, Inches(0.75), Inches(0.8), Inches(8.5), Inches(0.6),
         "SCHISM BMI Implementation", size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0.75), Inches(1.5), Inches(8.5), Inches(0.5),
         "Deep Dive: Architecture, Variables, NOAA NextGen, & Production Readiness",
         size=18, color=TITLE_ACCENT, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(2.5), Inches(7), Inches(0.4),
         "LynkerIntel/SCHISM_BMI Repository Analysis", size=16, color=LIGHT_BG, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.3), Inches(7), Inches(0.4),
         "Mohansai Sathram", size=20, color=LIGHT_BG, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.9), Inches(7), Inches(0.4),
         "Dr. Yao Hu", size=24, color=LIGHT_BG, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(4.6), Inches(7), Inches(0.4),
         "February 2026", size=16, color=LIGHT_BG, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(5.585), SLIDE_W, Inches(0.04), TITLE_ACCENT)


# ──────────────────────────────────────
# SLIDE 2: AGENDA / TABLE OF CONTENTS
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Presentation Overview")

topics = [
    ("1", "Repository &\nArchitecture", "Who built it, the two-layer\ndesign, and wrapper structure",
     BLUE_ACCENT, BLUE_LIGHT, "Slides 3\u20136"),
    ("2", "Variables &\nGrids", "17 exposed variables, 9 grids,\nand the t0/t1 data pattern",
     GREEN_ACCENT, GREEN_LIGHT, "Slides 7\u201311"),
    ("3", "NOAA Operations\n& NextGen", "ESMF production systems,\nNextGen experimental status,\n2024 coupling paper",
     ORANGE_ACCENT, ORANGE_LIGHT, "Slides 12\u201314"),
    ("4", "Production\nReadiness", "Validation status, limitations,\nand assessment for our project",
     PURPLE_ACCENT, PURPLE_LIGHT, "Slides 15\u201317"),
]

for i, (num, title, desc, accent, bg_color, slides) in enumerate(topics):
    x = Inches(0.25 + i * 2.45)
    y = Inches(1.2)

    add_shape(slide, x, y, Inches(2.25), Inches(3.8), bg_color, accent, Pt(1.5))

    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), y + Inches(0.2),
                                   Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.12), y + Inches(0.9), Inches(2.0), Inches(0.7),
             title, size=14, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.12), y + Inches(1.7), Inches(2.0), Inches(1.4),
             desc, size=10, color=MED_TEXT, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.12), y + Inches(3.3), Inches(2.0), Inches(0.3),
             slides, size=10, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────
# SLIDE 3: WHO BUILT IT & TIMELINE
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Who Built the SCHISM BMI and When?",
             "Built by Lynker Technologies (NOAA contractor) for the NextGen Water Resources Modeling Framework")

# People table
add_text(slide, Inches(0.5), Inches(1.15), Inches(9), Inches(0.25),
         "Key People:", size=13, color=DARK_TEXT, bold=True)

add_table(slide, Inches(0.4), Inches(1.45), Inches(9.2), 4, 3, [
    ["Person", "Role", "Organization"],
    ["Jason Ducker", "Primary developer, NextGen Forcings Director", "Lynker Technologies (NOAA OWP contractor)"],
    ["Joseph Zhang", "SCHISM lead developer, collaborator", "VIMS (Virginia Institute of Marine Science)"],
    ["Phil Miller", "NextGen framework integration lead", "NOAA OWP (Office of Water Prediction)"],
], col_widths=[Inches(2.0), Inches(4.0), Inches(3.2)])

# Timeline
add_text(slide, Inches(0.5), Inches(3.0), Inches(9), Inches(0.25),
         "Development Timeline:", size=13, color=DARK_TEXT, bold=True)

add_shape(slide, Inches(0.4), Inches(3.3), Inches(9.2), Inches(2.0), WHITE, GRAY_BORDER)

timeline = [
    ("Dec 2022", "Jason Ducker creates LynkerIntel/SCHISM_BMI repository", BLUE_ACCENT),
    ("Apr 2023", "schism-dev/schism_NWM_BMI created; Joseph Zhang collaborates on #ifdef design", GREEN_ACCENT),
    ("Jun 2023", "Issue #547 opened on NOAA-OWP/ngen: \"Evaluating SCHISM BMI as NextGen Formulation\"", ORANGE_ACCENT),
    ("Oct 2023", "Build documentation and wiki published", TEAL_ACCENT),
    ("Sep 2024", "Two-way NWM\u2194SCHISM coupling paper published (using ESMF, validates need for BMI)", PURPLE_ACCENT),
    ("Feb 2026", "Still experimental \u2014 not yet in production within NextGen framework", RED_ACCENT),
]

for i, (date, desc, accent) in enumerate(timeline):
    y_off = Inches(3.4 + i * 0.3)
    add_text(slide, Inches(0.6), y_off, Inches(1.2), Inches(0.25),
             date, size=9, color=accent, bold=True)
    add_text(slide, Inches(1.8), y_off, Inches(7.6), Inches(0.25),
             desc, size=9, color=MED_TEXT)


# ──────────────────────────────────────
# SLIDE 4: TWO-LAYER ARCHITECTURE
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "The Two Layers of SCHISM BMI",
             "SCHISM's BMI has two distinct layers that work together \u2014 inline #ifdef blocks + a separate wrapper file")

# Layer 1
add_shape(slide, Inches(0.4), Inches(1.15), Inches(4.5), Inches(2.2), BLUE_LIGHT, BLUE_ACCENT, Pt(1.5))
add_text(slide, Inches(0.6), Inches(1.2), Inches(4.1), Inches(0.3),
         "Layer 1: #ifdef Blocks (Inside SCHISM)", size=13, color=BLUE_ACCENT, bold=True)

add_bullets(slide, Inches(0.6), Inches(1.55), Inches(4.1), Inches(1.7), [
    "5 small code blocks in 3 source files:\n  schism_init.F90, schism_step.F90, misc_subs.F90",
    "Activated by: cmake -DUSE_NWM_BMI=ON",
    "Purpose: replace file-based forcing with\n  array-based data that BMI can write into",
    "When flag is OFF: blocks are skipped entirely,\n  SCHISM runs as a normal standalone model",
], size=10, bullet_color=BLUE_ACCENT, spacing=Pt(4))

# Layer 2
add_shape(slide, Inches(5.2), Inches(1.15), Inches(4.4), Inches(2.2), GREEN_LIGHT, GREEN_ACCENT, Pt(1.5))
add_text(slide, Inches(5.4), Inches(1.2), Inches(4.0), Inches(0.3),
         "Layer 2: BMI Wrapper (bmischism.f90)", size=13, color=GREEN_ACCENT, bold=True)

add_bullets(slide, Inches(5.4), Inches(1.55), Inches(4.0), Inches(1.7), [
    "Separate file: 1,729 lines of Fortran 2003",
    "Implements all 41 BMI functions by\n  extending the abstract bmi type from bmif_2_0",
    "Exposes 17 variables (12 input + 5 output)\n  across 9 different grids",
    "From LynkerIntel/SCHISM_BMI repo\n  (branch: bmi-integration-master)",
], size=10, bullet_color=GREEN_ACCENT, spacing=Pt(4))

# How they work together
add_code(slide, Inches(0.4), Inches(3.55), Inches(9.2), Inches(1.8),
"""    External Caller (NextGen / PyMT / Python)
              |
              |  BMI calls: initialize(), update(), get_value(), set_value()
              v
    +---------------------------+
    |   bmischism.f90            |  <-- Layer 2: Standard BMI interface (1,729 lines)
    |   (BMI Wrapper)            |      Maps 17 variables, 9 grids, IRF control
    +-------------+-------------+
                  |
                  |  Internal calls: schism_init(), schism_step(), schism_finalize()
                  |  Direct array access: ath3, eta2, pr1/pr2, windx1/windx2, etc.
                  v
    +---------------------------+
    |   SCHISM Source Code       |  <-- Layer 1: #ifdef blocks make arrays
    |   (with 5 #ifdef blocks)  |      writable from outside instead of files
    +---------------------------+""", size=8)


# ──────────────────────────────────────
# SLIDE 5: REPOSITORY STRUCTURE
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "LynkerIntel/SCHISM_BMI Repository",
             "GitHub: github.com/LynkerIntel/SCHISM_BMI \u2014 Branch: bmi-integration-master \u2014 Added as git submodule in our project")

# Left: Key files
add_shape(slide, Inches(0.4), Inches(1.15), Inches(4.5), Inches(2.6), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.6), Inches(1.2), Inches(4.1), Inches(0.3),
         "Key Files in src/BMI/", size=13, color=BLUE_ACCENT, bold=True)

add_table(slide, Inches(0.5), Inches(1.55), Inches(4.3), 6, 3, [
    ["File", "Lines", "Purpose"],
    ["bmischism.f90", "1,729", "Full BMI wrapper (all 41 functions)"],
    ["schism_bmi_driver_test.f90", "790", "Test driver (complete lifecycle test)"],
    ["bmi.f90", "564", "BMI abstract interface (spec copy)"],
    ["schism_model_container.f90", "~50", "Config + timestep state container"],
    ["LIMITATIONS", "10", "Official known limitations file"],
], col_widths=[Inches(2.0), Inches(0.5), Inches(1.8)])

# Right: Build config
add_shape(slide, Inches(5.2), Inches(1.15), Inches(4.4), Inches(2.6), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.4), Inches(1.2), Inches(4.0), Inches(0.3),
         "Build Configuration", size=13, color=GREEN_ACCENT, bold=True)

add_text(slide, Inches(5.4), Inches(1.55), Inches(4.0), Inches(0.2),
         "SCHISM.local.bmi_serial (cmake settings):", size=10, color=LIGHT_TEXT)

add_code(slide, Inches(5.3), Inches(1.8), Inches(4.2), Inches(1.3),
"""# Key cmake flags for BMI build:
set(USE_BMI ON)           # Enable BMI wrapper
set(OLDIO ON)             # Required for serial mode
set(BUILD_SHARED_LIBS ON) # Produce .so library
set(NO_PARMETIS ON)       # No mesh partitioning
set(TVD_LIM VL)           # Numerics setting""", size=9)

add_bullets(slide, Inches(5.4), Inches(3.2), Inches(4.0), Inches(0.5), [
    "Produces: libschism_bmi.so (shared library)",
], size=10, bullet_color=GREEN_ACCENT, spacing=Pt(3))

# Bottom: Our setup
add_shape(slide, Inches(0.4), Inches(3.95), Inches(9.2), Inches(1.3), TEAL_LIGHT, TEAL_ACCENT, Pt(1))
add_text(slide, Inches(0.6), Inches(4.0), Inches(8.8), Inches(0.25),
         "How We Have It Set Up in Our Project:", size=12, color=TEAL_ACCENT, bold=True)

add_code(slide, Inches(0.5), Inches(4.3), Inches(8.9), Inches(0.85),
"""# Added as git submodule pointing to bmi-integration-master branch:
git submodule add -b bmi-integration-master https://github.com/LynkerIntel/SCHISM_BMI.git SCHISM_BMI

# Our .gitmodules now has 5 submodules:
# wrf_hydro_nwm_public, bmi-fortran, bmi-example-fortran, schism_NWM_BMI, SCHISM_BMI""", size=9)


# ──────────────────────────────────────
# SLIDE 6: WRAPPER DESIGN PATTERNS
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "BMI Wrapper Design: bmischism.f90 (1,729 Lines)",
             "Follows the same patterns as the BMI Heat example \u2014 type extension, model container, select case dispatch")

# Pattern 1: Type Extension
add_shape(slide, Inches(0.4), Inches(1.15), Inches(4.5), Inches(1.7), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.55), Inches(1.2), Inches(4.2), Inches(0.25),
         "Pattern 1: Type Extension + Model Container", size=11, color=BLUE_ACCENT, bold=True)

add_code(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(1.25),
"""type, extends(bmi) :: bmi_schism
  private
  type(schism_model) :: model  ! Container holds:
  ! - iths (current timestep counter)
  ! - ntime (total timesteps)
  ! - time_step_size (seconds)
  ! - given_communicator (MPI handle)
contains
  procedure :: initialize => schism_initialize
  procedure :: update => schism_update  ! ...53 total
end type""", size=8)

# Pattern 2: IRF mapping
add_shape(slide, Inches(5.2), Inches(1.15), Inches(4.4), Inches(1.7), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.35), Inches(1.2), Inches(4.1), Inches(0.25),
         "Pattern 2: IRF Mapping to SCHISM Calls", size=11, color=GREEN_ACCENT, bold=True)

add_code(slide, Inches(5.3), Inches(1.5), Inches(4.2), Inches(1.25),
"""! initialize() -> calls schism_init()
!   Reads mesh, allocates arrays, sets BCs

! update() -> increments iths, then:
!   calls schism_step(iths)
!   (one ocean timestep: solve hydro eqns)

! finalize() -> calls schism_finalize()
!   Deallocates arrays, closes files""", size=8)

# Pattern 3: Select Case for variables
add_shape(slide, Inches(0.4), Inches(3.05), Inches(4.5), Inches(2.2), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.55), Inches(3.1), Inches(4.2), Inches(0.25),
         "Pattern 3: Variable Dispatch (select case)", size=11, color=ORANGE_ACCENT, bold=True)

add_code(slide, Inches(0.5), Inches(3.4), Inches(4.3), Inches(1.75),
"""! get_value("ETA2", output_array):
select case(name)
  case("ETA2")
    dest = eta2(1:np)  ! Copy water elevation
    bmi_status = BMI_SUCCESS
  case("VX")
    dest = uu2(1:np, nvrt) ! Surface velocity
    bmi_status = BMI_SUCCESS
  case default
    bmi_status = BMI_FAILURE ! Unknown var
end select""", size=8)

# Pattern 4: Global state access
add_shape(slide, Inches(5.2), Inches(3.05), Inches(4.4), Inches(2.2), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.35), Inches(3.1), Inches(4.1), Inches(0.25),
         "Pattern 4: Global State Access (not embedded)", size=11, color=PURPLE_ACCENT, bold=True)

add_bullets(slide, Inches(5.4), Inches(3.4), Inches(4.0), Inches(1.8), [
    "Unlike Heat model (embeds state in type),\n  SCHISM stores state in global module variables\n  (schism_glbl.F90 \u2014 eta2, uu2, vv2, ath3, etc.)",
    "Wrapper accesses globals via \"use schism_glbl\"\n  statements \u2014 reads/writes arrays directly",
    "The model container only holds config +\n  timestep tracking (not physics state)",
    "This is what we'll do for WRF-Hydro too \u2014\n  access RT_DOMAIN(1)%QLINK via module globals",
], size=9, bullet_color=PURPLE_ACCENT, spacing=Pt(4))


# ──────────────────────────────────────
# SLIDE 7: INPUT VARIABLES (12)
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Input Variables: What Goes INTO SCHISM via BMI (12 Variables)",
             "These are values that external models (like WRF-Hydro) or frameworks (NextGen) SEND to SCHISM")

# Boundary condition vars
add_shape(slide, Inches(0.4), Inches(1.15), Inches(9.2), Inches(0.35), BLUE_LIGHT, BLUE_ACCENT)
add_text(slide, Inches(0.6), Inches(1.17), Inches(8.8), Inches(0.3),
         "Boundary Condition Variables (3) \u2014 Water entering/leaving the ocean domain",
         size=11, color=BLUE_ACCENT, bold=True)

add_table(slide, Inches(0.4), Inches(1.55), Inches(9.2), 4, 4, [
    ["Variable Name", "What It Is", "Units", "SCHISM Internal Array"],
    ["Q_bnd_source", "River discharge at source elements (WRF-Hydro sends this!)", "m\u00B3/s", "ath3 (volume source)"],
    ["Q_bnd_sink", "Water withdrawal at sink elements", "m\u00B3/s", "ath3 (volume sink)"],
    ["ETA2_bnd", "Water level at open ocean boundaries (tides/surge)", "m", "ath2 (boundary condition)"],
], col_widths=[Inches(1.5), Inches(4.0), Inches(0.8), Inches(2.9)],
   header_bg=BLUE_ACCENT)

# Atmospheric forcing vars
add_shape(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(0.35), GREEN_LIGHT, GREEN_ACCENT)
add_text(slide, Inches(0.6), Inches(2.92), Inches(8.8), Inches(0.3),
         "Atmospheric Forcing Variables (6) \u2014 Weather data driving the ocean model",
         size=11, color=GREEN_ACCENT, bold=True)

add_table(slide, Inches(0.4), Inches(3.3), Inches(9.2), 7, 4, [
    ["Variable Name", "What It Is", "Units", "SCHISM Internal Array"],
    ["SFCPRS", "Surface atmospheric pressure", "Pa", "pr1/pr2 (two timepoints)"],
    ["TMP2m", "Air temperature at 2m above surface", "K", "airt1/airt2"],
    ["U10m", "Eastward wind speed at 10m", "m/s", "windx1/windx2"],
    ["V10m", "Northward wind speed at 10m", "m/s", "windy1/windy2"],
    ["SPFH2m", "Specific humidity at 2m", "kg/kg", "shum1/shum2"],
    ["RAINRATE", "Precipitation rate (converted to discharge!)", "kg/m\u00B2/s", "Added to ath3 source"],
], col_widths=[Inches(1.5), Inches(4.0), Inches(0.8), Inches(2.9)],
   header_bg=GREEN_ACCENT)

# Control vars (small)
add_shape(slide, Inches(0.4), Inches(5.0), Inches(9.2), Inches(0.35), ORANGE_LIGHT, ORANGE_ACCENT)
add_text(slide, Inches(0.6), Inches(5.02), Inches(2.0), Inches(0.3),
         "Control Variables (3):", size=10, color=ORANGE_ACCENT, bold=True)
add_text(slide, Inches(2.6), Inches(5.02), Inches(6.8), Inches(0.3),
         "ETA2_dt (boundary update interval), Q_dt (discharge update interval), bmi_mpi_comm_handle (MPI communicator)",
         size=10, color=MED_TEXT)


# ──────────────────────────────────────
# SLIDE 8: OUTPUT VARIABLES (5)
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Output Variables: What Comes OUT of SCHISM via BMI (5 Variables)",
             "These are values that external models can READ from SCHISM after each timestep")

# Output table
add_table(slide, Inches(0.4), Inches(1.2), Inches(9.2), 6, 4, [
    ["Variable Name", "What It Is", "Units", "Source Inside SCHISM"],
    ["ETA2", "Sea surface water elevation at ALL mesh nodes", "m", "Hydrodynamic solver (main output!)"],
    ["VX", "Eastward water current velocity (surface layer)", "m/s", "Momentum equation solver"],
    ["VY", "Northward water current velocity (surface layer)", "m/s", "Momentum equation solver"],
    ["TROUTE_ETA2", "Water elevation at monitoring stations only", "m", "Extracted from eta2 at station.in locations"],
    ["BEDLEVEL", "Ocean floor elevation (bathymetry, static)", "m", "Depth array (inverted: positive = above datum)"],
], col_widths=[Inches(1.5), Inches(3.6), Inches(0.7), Inches(3.4)])

# The coupling pair
add_text(slide, Inches(0.5), Inches(3.4), Inches(9), Inches(0.25),
         "The Critical Coupling Pair (WRF-Hydro \u2194 SCHISM):", size=13, color=DARK_TEXT, bold=True)

add_code(slide, Inches(0.4), Inches(3.7), Inches(9.2), Inches(1.3),
"""   WRF-HYDRO (our BMI wrapper)                              SCHISM (existing BMI)
   ============================                              ====================

   QLINK (streamflow)  -----> Q_bnd_source (BMI input) ----> ath3 (source discharge)
   CSDMS: channel_water__volume_flow_rate                     River water enters ocean

   (future: receive     <---- ETA2 (BMI output)        <---- eta2 (water elevation)
    coastal water level)                                      Coastal surge pushes upstream""", size=9)

# Summary box
add_shape(slide, Inches(0.4), Inches(5.1), Inches(9.2), Inches(0.35), BLUE_LIGHT, BLUE_ACCENT)
add_text(slide, Inches(0.6), Inches(5.12), Inches(8.8), Inches(0.3),
         "Total: 12 input + 5 output = 17 variables  |  Compare: Our WRF-Hydro plan = 2 input + 8 output = 10 variables  |  Heat example = 1 variable",
         size=10, color=BLUE_ACCENT, bold=True)


# ──────────────────────────────────────
# SLIDE 9: THE 9-GRID SYSTEM
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "The 9-Grid System: How SCHISM Describes Its Mesh",
             "Different variables live on different parts of the unstructured triangular mesh \u2014 BMI needs a separate grid for each")

# Spatial grids (6)
add_text(slide, Inches(0.5), Inches(1.15), Inches(4.5), Inches(0.25),
         "Spatial Grids (6):", size=12, color=BLUE_ACCENT, bold=True)

add_table(slide, Inches(0.4), Inches(1.45), Inches(5.0), 7, 4, [
    ["Grid", "Name", "Type", "Variables"],
    ["1", "All Nodes", "unstructured", "ETA2, VX, VY, BEDLEVEL"],
    ["2", "All Elements", "points", "RAINRATE"],
    ["3", "Offshore Boundary", "points", "ETA2_bnd"],
    ["4", "Source Elements", "points", "Q_bnd_source"],
    ["5", "Sink Elements", "points", "Q_bnd_sink"],
    ["6", "Station Points", "points", "TROUTE_ETA2"],
], col_widths=[Inches(0.5), Inches(1.5), Inches(1.1), Inches(1.9)])

# Virtual scalar grids (3)
add_text(slide, Inches(5.6), Inches(1.15), Inches(4.0), Inches(0.25),
         "Virtual Scalar Grids (3):", size=12, color=GREEN_ACCENT, bold=True)

add_table(slide, Inches(5.5), Inches(1.45), Inches(4.1), 4, 3, [
    ["Grid", "Name", "Variable"],
    ["7", "ETA2 Timestep", "ETA2_dt"],
    ["8", "Q Timestep", "Q_dt"],
    ["9", "MPI Communicator", "bmi_mpi_comm_handle"],
], col_widths=[Inches(0.5), Inches(1.6), Inches(2.0)],
   header_bg=GREEN_ACCENT)

# Visual diagram
add_code(slide, Inches(0.4), Inches(3.65), Inches(5.0), Inches(1.7),
"""  SCHISM Ocean Domain (Triangular Mesh)
  +---+---+---+---+---+---+---+
  | / | \\ | / | \\ | / | \\ | / |  Grid 1: All Nodes
  +---+---+---+---+---+---+---+    (ETA2, VX, VY, BEDLEVEL)
  | \\ | / | \\ | / | \\ | / | \\ |
  +---+---+---+---+---+---+---+  Grid 2: Element centers
  ^                           ^    (RAINRATE -> discharge)
  |                           |
  Grid 3: Ocean boundary      Grid 4: Source elements
  (ETA2_bnd)                  (Q_bnd_source = river mouths)""", size=8)

# Comparison to WRF-Hydro
add_shape(slide, Inches(5.7), Inches(3.65), Inches(3.9), Inches(1.7), ORANGE_LIGHT, ORANGE_ACCENT, Pt(1))
add_text(slide, Inches(5.85), Inches(3.7), Inches(3.6), Inches(0.25),
         "Comparison to WRF-Hydro:", size=11, color=ORANGE_ACCENT, bold=True)

add_table(slide, Inches(5.8), Inches(4.0), Inches(3.7), 4, 2, [
    ["WRF-Hydro Grid", "SCHISM Grid"],
    ["Grid 0: 1km LSM", "Grid 1: All Nodes (unstructured)"],
    ["Grid 1: 250m routing", "Grid 2: All Elements"],
    ["Grid 2: Channel network", "Grids 3-6: Subsets of mesh"],
], col_widths=[Inches(1.8), Inches(1.9)],
   header_bg=ORANGE_ACCENT)


# ──────────────────────────────────────
# SLIDE 10: THE t0/t1 TWO-TIMEPOINT PATTERN
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "The Two-Timepoint (t0/t1) Sliding Pattern",
             "SCHISM stores forcing data at TWO time levels and interpolates between them \u2014 a clever design for coarse-to-fine coupling")

# Left: The problem
add_shape(slide, Inches(0.4), Inches(1.15), Inches(4.5), Inches(1.6), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.6), Inches(1.2), Inches(4.1), Inches(0.25),
         "The Problem:", size=12, color=RED_ACCENT, bold=True)

add_bullets(slide, Inches(0.6), Inches(1.5), Inches(4.1), Inches(1.2), [
    "SCHISM runs with fine timesteps (~100 sec)",
    "External data arrives at coarse intervals (~1 hr)",
    "Without interpolation, forcing \"jumps\" abruptly\n  between updates \u2014 causing numerical instability",
], size=10, bullet_color=RED_ACCENT, spacing=Pt(4))

# Right: The solution
add_shape(slide, Inches(5.2), Inches(1.15), Inches(4.4), Inches(1.6), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.4), Inches(1.2), Inches(4.0), Inches(0.25),
         "The Solution: t0/t1 Sliding", size=12, color=GREEN_ACCENT, bold=True)

add_bullets(slide, Inches(5.4), Inches(1.5), Inches(4.0), Inches(1.2), [
    "Every forcing variable stored as TWO snapshots:\n  t0 = old value, t1 = new value",
    "SCHISM interpolates between t0 and t1 during\n  each internal timestep (smooth transition)",
    "Each set_value() call slides: old t1 \u2192 new t0,\n  new data \u2192 new t1 (automatic bookkeeping)",
], size=10, bullet_color=GREEN_ACCENT, spacing=Pt(4))

# Diagram
add_code(slide, Inches(0.4), Inches(2.95), Inches(9.2), Inches(2.35),
"""BEFORE set_value("Q_bnd_source", new_discharge):

   t0 = [100, 200, 150]    <-- old values (from 2 updates ago)
   t1 = [120, 180, 160]    <-- current values (from last update)

   SCHISM interpolates between t0 and t1 for each internal timestep

AFTER set_value("Q_bnd_source", new_discharge):

   t0 = [120, 180, 160]    <-- old t1 SLIDES to become new t0
   t1 = [140, 220, 170]    <-- new data becomes new t1

   Now SCHISM interpolates between NEW t0 and t1

RESULT: Smooth forcing transitions, no jumps, caller only provides data every ~1 hour""", size=9)


# ──────────────────────────────────────
# SLIDE 11: RAINRATE SPECIAL CASE
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Special Case: How Rainfall Becomes River Discharge",
             "RAINRATE is the ONLY input variable that ADDS to existing values \u2014 all others REPLACE")

# The conversion
add_shape(slide, Inches(0.4), Inches(1.15), Inches(4.5), Inches(2.3), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.6), Inches(1.2), Inches(4.1), Inches(0.25),
         "The Conversion:", size=13, color=BLUE_ACCENT, bold=True)

add_code(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(1.85),
"""When set_value("RAINRATE", rainfall_array):

  For each mesh element i:
    extra_discharge(i) =
      rainfall_rate(i)  x  element_area(i)
              |                    |
         (kg/m2/s)              (m2)
              = m3/s (volume flow rate)

  Then ADD to existing source discharge:
    ath3(i) = ath3(i) + extra_discharge(i)

  (not replace -- accumulate!)""", size=9)

# Why this design
add_shape(slide, Inches(5.2), Inches(1.15), Inches(4.4), Inches(2.3), WHITE, GRAY_BORDER)
add_text(slide, Inches(5.4), Inches(1.2), Inches(4.0), Inches(0.25),
         "Why This Design?", size=13, color=GREEN_ACCENT, bold=True)

add_bullets(slide, Inches(5.4), Inches(1.55), Inches(4.0), Inches(1.8), [
    "SCHISM treats ALL water inputs as volume\n  sources (m\u00B3/s) \u2014 unified treatment",
    "Whether from a river or from rain falling on\n  the ocean, it's all \"water entering the system\"",
    "By converting to discharge and adding to ath3,\n  SCHISM handles inputs uniformly",
    "Different from WRF-Hydro, where rainfall is a\n  separate atmospheric forcing variable",
], size=10, bullet_color=GREEN_ACCENT, spacing=Pt(4))

# Warning banner
add_shape(slide, Inches(0.4), Inches(3.65), Inches(9.2), Inches(0.65), RED_LIGHT, RED_ACCENT, Pt(1.5))
add_text(slide, Inches(0.6), Inches(3.68), Inches(8.8), Inches(0.6),
         "IMPORTANT: RAINRATE is the ONLY variable that accumulates. All other set_value calls REPLACE existing\n"
         "values (with the t0/t1 slide pattern). This inconsistency can cause bugs if callers don't account for it.",
         size=11, color=RED_ACCENT, bold=True)

# The 6 variable pairs with t0/t1
add_text(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(0.25),
         "The t0/t1 Array Pairs for Atmospheric Forcing:", size=12, color=DARK_TEXT, bold=True)

add_table(slide, Inches(0.4), Inches(4.8), Inches(9.2), 4, 5, [
    ["BMI Variable", "t0 Array (old)", "t1 Array (new)", "Behavior", "Grid"],
    ["SFCPRS", "pr1", "pr2", "Replace (slide)", "All Nodes (#1)"],
    ["U10m / V10m", "windx1 / windy1", "windx2 / windy2", "Replace (slide)", "All Nodes (#1)"],
    ["RAINRATE", "N/A", "N/A", "Accumulate (add to ath3)", "All Elements (#2)"],
], col_widths=[Inches(1.3), Inches(1.8), Inches(1.8), Inches(2.2), Inches(2.1)])


# ──────────────────────────────────────
# SLIDE 12: NOAA USING SCHISM TODAY (ESMF)
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Where Is NOAA Using SCHISM Today?",
             "SCHISM runs operationally in NOAA forecasting systems \u2014 but using ESMF coupling, NOT BMI")

# Top banner
add_shape(slide, Inches(0.4), Inches(1.15), Inches(9.2), Inches(0.55), ORANGE_LIGHT, ORANGE_ACCENT, Pt(1.5))
add_text(slide, Inches(0.6), Inches(1.18), Inches(8.8), Inches(0.5),
         "Key Distinction: SCHISM in NOAA production uses ESMF/NUOPC coupling framework \u2014 NOT BMI.\n"
         "The BMI version is experimental, being developed for the NextGen framework.",
         size=11, color=ORANGE_ACCENT, bold=True)

# Operational systems table
add_text(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(0.25),
         "SCHISM in NOAA Operational Systems (ESMF-based, NOT BMI):", size=12, color=DARK_TEXT, bold=True)

add_table(slide, Inches(0.4), Inches(2.2), Inches(9.2), 5, 4, [
    ["System", "What It Does", "Status", "Coupling"],
    ["STOFS-3D-Atlantic", "Total water level forecasts for Atlantic/Gulf coasts\n(2.9M nodes, 5.7M elements)", "Operational since Jan 2023", "ESMF/NUOPC"],
    ["STOFS-3D-Pacific", "Total water level forecasts for Pacific coast", "Planned for late 2025", "ESMF/NUOPC"],
    ["NWM v3.0 TWL", "Total water level guidance combining river + ocean", "Operational since summer 2023", "ESMF/NUOPC"],
    ["ICOGS", "Inland-Coastal Operational Guidance System", "Pre-operational since April 2021", "Direct coupling"],
], col_widths=[Inches(1.8), Inches(3.6), Inches(2.0), Inches(1.8)])

# Visual: current vs future
add_code(slide, Inches(0.4), Inches(4.0), Inches(9.2), Inches(1.3),
"""  TODAY (Production - working now):

  NWM ---- ESMF/NUOPC ----> SCHISM ----> STOFS forecasts (operational)
           (Earth System Modeling Framework)

  ================================================================

  FUTURE (In Development - what we're building toward):

  WRF-Hydro/NWM ---- BMI 2.0 ----> SCHISM ----> NextGen forecasts (experimental)
                     (Basic Model Interface)""", size=9)


# ──────────────────────────────────────
# SLIDE 13: NEXTGEN FRAMEWORK
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "The NextGen Framework: Where SCHISM BMI Is Headed",
             "NextGen = NOAA's next-generation modular water modeling framework. Any BMI model can be plugged in or swapped out.")

# Architecture diagram
add_code(slide, Inches(0.4), Inches(1.15), Inches(9.2), Inches(2.3),
"""                           NextGen Framework (ngen)
    +----------------------------------------------------------------------+
    |                                                                      |
    |  +----------------+   +----------------+   +------------------+      |
    |  | Atmospheric    |   | Hydrologic     |   | Channel Routing  |      |
    |  | Forcings (BMI) |   | Models (BMI)   |   | t-route (BMI)   |      |
    |  |                |   |                |   |                  |      |
    |  | * NWM Forcing  |   | * CFE          |   |  Discharge at   |      |
    |  | * ERA5         |   | * TOPMODEL     |   |  coastal outlets |      |
    |  | * AORC         |   | * NoahOWP      |   |                  |      |
    |  +-------+--------+   +-------+--------+   +--------+---------+      |
    |          |                    |                      |                |
    |          v                    v                      v                |
    |  +--------------------------------------------------------------+    |
    |  |           COASTAL REALIZATION (new! - in development)         |    |
    |  |   +-------------+              +-------------+                |    |
    |  |   | SCHISM (BMI) |     OR     | D-Flow FM    |                |    |
    |  |   | Water levels |            | (BMI)        |                |    |
    |  |   | Currents     |            |              |                |    |
    |  |   +-------------+              +-------------+                |    |
    |  +--------------------------------------------------------------+    |
    +----------------------------------------------------------------------+""", size=7)

# Status table
add_text(slide, Inches(0.5), Inches(3.6), Inches(9), Inches(0.25),
         "SCHISM in NextGen \u2014 Remaining Work Items:", size=12, color=DARK_TEXT, bold=True)

add_table(slide, Inches(0.4), Inches(3.9), Inches(9.2), 6, 3, [
    ["Task", "Status", "Notes"],
    ["BMI coastal formulation in ngen", "Completed", "SCHISM compiles as ngen formulation"],
    ["Lake Champlain test case", "Created", "Used for framework validation testing"],
    ["Define coastal realization config", "In progress", "Config structure being designed"],
    ["MPI process mapping for coastal models", "In progress", "Multiple cores for SCHISM"],
    ["Integration with t-route discharge", "In progress", "Map nexus IDs to BMI variable indices"],
], col_widths=[Inches(3.2), Inches(1.5), Inches(4.5)])

# Key quote
add_shape(slide, Inches(0.4), Inches(5.1), Inches(9.2), Inches(0.35), PURPLE_LIGHT, PURPLE_ACCENT)
add_text(slide, Inches(0.6), Inches(5.12), Inches(8.8), Inches(0.3),
         "GitHub Issue #547 (NOAA-OWP/ngen): \"Evaluating SCHISM BMI as a NextGen Formulation\" \u2014 Opened Jun 2023, still OPEN",
         size=10, color=PURPLE_ACCENT, bold=True)


# ──────────────────────────────────────
# SLIDE 14: 2024 TWO-WAY COUPLING PAPER
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "The 2024 Two-Way Coupling Paper: Key Findings",
             "Zhang et al., Hydrology (MDPI), Sep 2024 \u2014 Validates the scientific need for our WRF-Hydro BMI project")

# Paper details
add_shape(slide, Inches(0.4), Inches(1.15), Inches(4.5), Inches(1.1), BLUE_LIGHT, BLUE_ACCENT)
add_text(slide, Inches(0.6), Inches(1.18), Inches(4.1), Inches(0.25),
         "Paper Details:", size=11, color=BLUE_ACCENT, bold=True)
add_multiline(slide, Inches(0.6), Inches(1.4), Inches(4.1), Inches(0.8),
              "Two-Way Coupling of NWM and SCHISM for\nEnhanced Coastal Discharge Predictions\nTest case: Hurricane Matthew (Oct 2016, SC)\nCoupling: CoastalApp/ESMF (not BMI)",
              size=9, color=MED_TEXT, line_spacing=Pt(13))

# Finding 1
add_shape(slide, Inches(5.2), Inches(1.15), Inches(4.4), Inches(1.1), GREEN_LIGHT, GREEN_ACCENT)
add_text(slide, Inches(5.4), Inches(1.18), Inches(4.0), Inches(0.25),
         "Finding 1: Two-Way Is Better", size=11, color=GREEN_ACCENT, bold=True)
add_multiline(slide, Inches(5.4), Inches(1.4), Inches(4.0), Inches(0.8),
              "One-way: NWM doesn't know about storm surge,\nrivers keep flowing normally during storms\nTwo-way: rivers slow when ocean pushes back,\nmuch better discharge predictions during storms",
              size=9, color=MED_TEXT, line_spacing=Pt(13))

# Finding 2
add_shape(slide, Inches(0.4), Inches(2.45), Inches(4.5), Inches(1.0), ORANGE_LIGHT, ORANGE_ACCENT)
add_text(slide, Inches(0.6), Inches(2.48), Inches(4.1), Inches(0.25),
         "Finding 2: Tidal Effects on Rivers", size=11, color=ORANGE_ACCENT, bold=True)
add_multiline(slide, Inches(0.6), Inches(2.7), Inches(4.1), Inches(0.7),
              "Daily tides cause river discharge to oscillate.\nOne-way NWM completely misses this.\nTwo-way captures tidal \"breathing\" of rivers.",
              size=9, color=MED_TEXT, line_spacing=Pt(13))

# Finding 3
add_shape(slide, Inches(5.2), Inches(2.45), Inches(4.4), Inches(1.0), PURPLE_LIGHT, PURPLE_ACCENT)
add_text(slide, Inches(5.4), Inches(2.48), Inches(4.0), Inches(0.25),
         "Finding 3: Peak Discharge Timing", size=11, color=PURPLE_ACCENT, bold=True)
add_multiline(slide, Inches(5.4), Inches(2.7), Inches(4.0), Inches(0.7),
              "Two-way more accurately predicts WHEN peak\ndischarge occurs (not just how much).\nAccounts for floodplain water storage.",
              size=9, color=MED_TEXT, line_spacing=Pt(13))

# Finding 4: Why this matters
add_shape(slide, Inches(0.4), Inches(3.65), Inches(9.2), Inches(0.65), TEAL_LIGHT, TEAL_ACCENT, Pt(1.5))
add_text(slide, Inches(0.6), Inches(3.68), Inches(8.8), Inches(0.6),
         "Finding 4: The paper explicitly calls for ESMF \u2192 BMI transition for NextGen. Quote: \"The proposed shift\n"
         "from ESMF to BMI for coastal coupling requires innovative solutions.\" THIS IS EXACTLY WHAT WE'RE BUILDING.",
         size=11, color=TEAL_ACCENT, bold=True)

# Why it matters for us
add_shape(slide, Inches(0.4), Inches(4.5), Inches(9.2), Inches(0.85), WHITE, GRAY_BORDER)
add_text(slide, Inches(0.6), Inches(4.55), Inches(8.8), Inches(0.25),
         "Why This Paper Validates Our Project:", size=12, color=DARK_TEXT, bold=True)

add_bullets(slide, Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.5), [
    "Two-way coupling (requires BMI on BOTH models) produces significantly better results",
    "ESMF \u2192 BMI transition is an active research priority recognized by NOAA",
    "The specific variables we're exposing (streamflow, water elevation) are exactly what the coupling needs",
], size=9, bullet_color=BLUE_ACCENT, spacing=Pt(2))


# ──────────────────────────────────────
# SLIDE 15: VALIDATION STATUS
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Does It Work? Current Validation Status",
             "The SCHISM BMI framework works, but the specific coupling variables we need are flagged as 'not fully validated'")

# What works
add_shape(slide, Inches(0.4), Inches(1.15), Inches(4.5), Inches(2.8), GREEN_LIGHT, GREEN_ACCENT, Pt(1))
add_text(slide, Inches(0.6), Inches(1.2), Inches(4.1), Inches(0.3),
         "What Works", size=14, color=GREEN_ACCENT, bold=True)

add_table(slide, Inches(0.5), Inches(1.55), Inches(4.3), 8, 2, [
    ["Aspect", "Status"],
    ["BMI interface compiles", "Working"],
    ["Initialize/Update/Finalize cycle", "Working"],
    ["Variable querying (names, types, units)", "Working"],
    ["Grid metadata (node counts, connectivity)", "Working"],
    ["MPI communicator passing", "Working"],
    ["Atmospheric forcing via set_value", "Working"],
    ["ETA2 output via get_value", "Working"],
], col_widths=[Inches(2.5), Inches(1.8)],
   header_bg=GREEN_ACCENT)

# What doesn't work
add_shape(slide, Inches(5.2), Inches(1.15), Inches(4.4), Inches(2.8), RED_LIGHT, RED_ACCENT, Pt(1))
add_text(slide, Inches(5.4), Inches(1.2), Inches(4.0), Inches(0.3),
         "NOT Fully Validated", size=14, color=RED_ACCENT, bold=True)

add_table(slide, Inches(5.3), Inches(1.55), Inches(4.2), 5, 2, [
    ["Aspect", "Issue"],
    ["Water level BCs (ETA2_bnd)", "Interpolation timing edge cases"],
    ["Source/sink (Q_bnd_source/sink)", "Discharge accumulation untested"],
    ["Two-way data exchange", "No real coupling test via BMI yet"],
    ["Multi-instance SCHISM", "Not tested for NextGen"],
], col_widths=[Inches(2.0), Inches(2.2)],
   header_bg=RED_ACCENT)

# Official LIMITATIONS file
add_shape(slide, Inches(0.4), Inches(4.15), Inches(9.2), Inches(1.15), ORANGE_LIGHT, ORANGE_ACCENT, Pt(1.5))
add_text(slide, Inches(0.6), Inches(4.2), Inches(8.8), Inches(0.25),
         "From the Official LIMITATIONS File (SCHISM_BMI/src/BMI/LIMITATIONS):", size=12, color=ORANGE_ACCENT, bold=True)

add_bullets(slide, Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.7), [
    "\"MPI parallelization is partially implemented for SCHISM driven through BMI\"",
    "\"Water level and source/sink connectors have not been fully validated yet for SCHISM\"",
], size=11, bullet_color=ORANGE_ACCENT, spacing=Pt(4))


# ──────────────────────────────────────
# SLIDE 16: KNOWN LIMITATIONS (DETAILED)
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Known Limitations from Code Analysis",
             "Beyond the official LIMITATIONS file, code analysis reveals additional constraints and design choices")

add_table(slide, Inches(0.4), Inches(1.15), Inches(9.2), 9, 3, [
    ["Limitation", "Details", "Impact"],
    ["No grid shape/spacing/origin", "Returns BMI_FAILURE \u2014 unstructured mesh\nhas no uniform spacing", "Grid functions partially unsupported"],
    ["Surface layer only for VX/VY", "Only top vertical layer of velocity exposed,\nnot the full 3D field", "Limited depth profile access"],
    ["RAINRATE accumulates", "Adds to ath3 instead of replacing \u2014\ninconsistent with other set_value calls", "Can cause bugs if not accounted for"],
    ["No set_value for outputs", "Cannot write INTO eta2 or VX/VY from\noutside \u2014 they are read-only", "One-way data flow for outputs"],
    ["First timestep cold start", "At t=0, no boundary data exists yet.\nSCHISM uses zeros until first set_value", "May cause initial transients"],
    ["OLDIO required for serial", "Must use OLDIO=ON for single-core runs.\nScribe I/O doesn't work in serial mode", "Build configuration constraint"],
    ["No CSDMS Standard Names", "Uses custom names (Q_bnd_source, ETA2)\nnot CSDMS (channel_water__volume_flow_rate)", "Not compatible with PyMT catalog"],
    ["Not registered with CSDMS", "No pymt_schism package exists.\nNot in CSDMS model repository", "Can't use PyMT coupling framework"],
], col_widths=[Inches(2.0), Inches(4.2), Inches(3.0)])


# ──────────────────────────────────────
# SLIDE 17: PRODUCTION READINESS ASSESSMENT
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Production Readiness Assessment",
             "Is the SCHISM BMI ready for production use? A clear-eyed evaluation of the current state.")

# Big verdict
add_shape(slide, Inches(0.4), Inches(1.15), Inches(9.2), Inches(0.7), ORANGE_LIGHT, ORANGE_ACCENT, Pt(2))
add_text(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.6),
         "Verdict: NOT production ready. The SCHISM BMI is a functional experimental prototype.\n"
         "The framework integration works, but key coupling variables remain unvalidated.",
         size=13, color=ORANGE_ACCENT, bold=True)

# Three assessment cards
cards = [
    ("Framework\nReadiness", "The BMI interface (all 41 functions)\ncompiles and runs correctly.\nInit/update/finalize cycle works.\nTest driver validates lifecycle.\n\nThis part IS production-ready.",
     GREEN_ACCENT, GREEN_LIGHT, "7/10"),
    ("Coupling\nReadiness", "Water level + source/sink (the\nexact variables needed for coupling)\nare NOT fully validated.\nNo real two-model coupling test\nhas been run through BMI yet.\n\nThis part is NOT ready.",
     RED_ACCENT, RED_LIGHT, "3/10"),
    ("NextGen\nIntegration", "SCHISM compiles as ngen formulation.\nLake Champlain test case exists.\nBut MPI mapping, config structure,\nand t-route integration are all\nstill in progress.\n\nThis part is IN PROGRESS.",
     ORANGE_ACCENT, ORANGE_LIGHT, "5/10"),
]

for i, (title, desc, accent, bg, score) in enumerate(cards):
    x = Inches(0.3 + i * 3.2)
    y = Inches(2.1)
    add_shape(slide, x, y, Inches(3.0), Inches(3.2), bg, accent, Pt(1.5))

    # Score badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.05), y + Inches(0.1),
                                   Inches(0.85), Inches(0.85))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = score
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.15), y + Inches(1.05), Inches(2.7), Inches(0.5),
             title, size=13, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    add_multiline(slide, x + Inches(0.15), y + Inches(1.55), Inches(2.7), Inches(1.5),
                  desc, size=8, color=MED_TEXT, line_spacing=Pt(12))


# ──────────────────────────────────────
# SLIDE 18: KEY TAKEAWAYS FOR OUR PROJECT
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Key Takeaways for Our WRF-Hydro BMI",
             "What we can learn from the SCHISM BMI implementation \u2014 patterns to follow and improvements to make")

# Patterns to follow
add_shape(slide, Inches(0.4), Inches(1.15), Inches(4.5), Inches(2.4), GREEN_LIGHT, GREEN_ACCENT, Pt(1))
add_text(slide, Inches(0.6), Inches(1.2), Inches(4.1), Inches(0.3),
         "Design Patterns to Follow:", size=13, color=GREEN_ACCENT, bold=True)

add_table(slide, Inches(0.5), Inches(1.55), Inches(4.3), 7, 2, [
    ["SCHISM Pattern", "Apply to WRF-Hydro"],
    ["Separate wrapper file", "Write bmi_wrf_hydro.f90"],
    ["Model container for config", "Container for WRF-Hydro state"],
    ["t0/t1 pattern for forcing", "Consider for input variables"],
    ["Multiple grids for domains", "Grid 0 (1km), Grid 1 (250m), Grid 2"],
    ["MPI communicator via set_value", "Same pattern for parallel"],
    ["Serial-first approach", "Start serial, add MPI later"],
], col_widths=[Inches(2.0), Inches(2.3)],
   header_bg=GREEN_ACCENT)

# Improvements over SCHISM
add_shape(slide, Inches(5.2), Inches(1.15), Inches(4.4), Inches(2.4), BLUE_LIGHT, BLUE_ACCENT, Pt(1))
add_text(slide, Inches(5.4), Inches(1.2), Inches(4.0), Inches(0.3),
         "Our Improvements Over SCHISM:", size=13, color=BLUE_ACCENT, bold=True)

add_table(slide, Inches(5.3), Inches(1.55), Inches(4.2), 5, 2, [
    ["SCHISM Approach", "Our Improvement"],
    ["Custom names (Q_bnd_source)", "CSDMS Standard Names"],
    ["Not registered with CSDMS", "Plan to register + babelize"],
    ["No pymt_schism plugin", "Create pymt_wrfhydro"],
    ["Limited validation", "Validate vs standalone from day 1"],
], col_widths=[Inches(2.0), Inches(2.2)],
   header_bg=BLUE_ACCENT)

# Variable count comparison
add_text(slide, Inches(0.5), Inches(3.75), Inches(9), Inches(0.25),
         "Variable Count Comparison:", size=12, color=DARK_TEXT, bold=True)

add_table(slide, Inches(0.4), Inches(4.05), Inches(9.2), 4, 5, [
    ["Model", "Input Vars", "Output Vars", "Total", "Grids"],
    ["SCHISM BMI (existing)", "12", "5", "17", "9"],
    ["WRF-Hydro BMI (our plan)", "2", "8", "10", "3"],
    ["BMI Heat Example (reference)", "0", "1", "1", "1"],
], col_widths=[Inches(2.8), Inches(1.3), Inches(1.3), Inches(1.3), Inches(2.5)])

# Bottom note
add_shape(slide, Inches(0.4), Inches(5.1), Inches(9.2), Inches(0.35), PURPLE_LIGHT, PURPLE_ACCENT)
add_text(slide, Inches(0.6), Inches(5.12), Inches(8.8), Inches(0.3),
         "Our WRF-Hydro BMI scope is between the simple heat example and the complex SCHISM \u2014 a very reasonable starting point.",
         size=10, color=PURPLE_ACCENT, bold=True)


# ──────────────────────────────────────
# SLIDE 19: SUMMARY / KEY FACTS
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)
slide_header(slide, "Summary: SCHISM BMI Key Facts")

add_code(slide, Inches(0.4), Inches(1.15), Inches(9.2), Inches(3.5),
"""                          SCHISM BMI -- Key Facts
  ============================================================================

  WHO:   Jason Ducker (Lynker/NOAA OWP) + Joseph Zhang (VIMS)
  WHEN:  Development started December 2022
  WHERE: LynkerIntel/SCHISM_BMI repo, branch: bmi-integration-master
  WHAT:  bmischism.f90 (1,729 lines) + 5 #ifdef blocks in SCHISM source

  VARIABLES: 12 input + 5 output = 17 total
  GRIDS:     9 (6 spatial + 3 virtual scalar)
  KEY IN:    Q_bnd_source (discharge from rivers -- what WRF-Hydro sends)
  KEY OUT:   ETA2 (water elevation -- for 2-way coupling back to WRF-Hydro)

  STATUS:
    [OK]  Framework works (init/update/finalize cycle)
    [!!]  Water level + source/sink NOT fully validated
    [NO]  Not in production -- experimental for NextGen
    [NO]  Not registered with CSDMS, no pymt_schism

  NOAA TODAY:     SCHISM runs in STOFS/NWM via ESMF (not BMI)
  NOAA FUTURE:    SCHISM in NextGen via BMI (in development)
  OUR PROJECT:    Building the WRF-Hydro side of this coupling""", size=9)

# Bottom insight
add_shape(slide, Inches(0.4), Inches(4.85), Inches(9.2), Inches(0.55), BLUE_LIGHT, BLUE_ACCENT, Pt(1.5))
add_text(slide, Inches(0.6), Inches(4.88), Inches(8.8), Inches(0.5),
         "Bottom Line: The SCHISM BMI is a solid reference implementation, but NOT production ready. Our WRF-Hydro BMI\n"
         "can learn from its patterns while improving with CSDMS Standard Names and proper validation.",
         size=11, color=BLUE_ACCENT, bold=True)


# ──────────────────────────────────────
# SLIDE 20: THANK YOU
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)

add_text(slide, Inches(1), Inches(1.2), Inches(8), Inches(0.6),
         "Thank You", size=36, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(2.0), Inches(8), Inches(0.4),
         "Questions & Discussion", size=20, color=BLUE_ACCENT, alignment=PP_ALIGN.CENTER)

# Summary cards
summary_cards = [
    ("Architecture", "Two-layer design:\n5 inline #ifdef blocks +\n1,729-line bmischism.f90\nwrapper. Follows same\npatterns as BMI Heat.", BLUE_ACCENT, BLUE_LIGHT),
    ("NOAA Usage", "SCHISM runs operationally\nin STOFS via ESMF.\nBMI version is experimental\nfor NextGen framework.\nNot production ready.", ORANGE_ACCENT, ORANGE_LIGHT),
    ("Our Project", "WRF-Hydro BMI wrapper\nlearns from SCHISM patterns.\n10 variables, 3 grids.\nCSDS names + validation\nfrom day one.", GREEN_ACCENT, GREEN_LIGHT),
]

for i, (title, desc, accent, bg) in enumerate(summary_cards):
    x = Inches(0.7 + i * 3.1)
    y = Inches(2.8)
    add_shape(slide, x, y, Inches(2.8), Inches(2.2), bg, accent, Pt(1.5))
    add_text(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.5), Inches(0.35),
             title, size=14, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.5), Inches(1.5),
                  desc, size=10, color=MED_TEXT, line_spacing=Pt(14))

# Bottom bar
add_rect(slide, Inches(0), Inches(5.3), SLIDE_W, Inches(0.325), BLUE_ACCENT)
add_text(slide, Inches(1), Inches(5.32), Inches(8), Inches(0.3),
         "SCHISM BMI Deep Dive  |  Mohansai Sathram  |  February 2026",
         size=10, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════
output_path = "/mnt/c/Users/mohansai/Desktop/Projects/VS_Code/WRF-Hydro-BMI/bmi_wrf_hydro/Docs/Weekly Reporting/SCHISM_BMI_Deep_Dive.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
